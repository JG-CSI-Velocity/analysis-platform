"""
deck_builder.py
---------------
PowerPoint automation for client report generation.

This module provides:
    - DECK_CONFIG: Global configuration dictionary for all settings
    - SlideContent: Dataclass to define what goes on each slide
    - DeckBuilder: Class that assembles slides into a PowerPoint deck
    - setup_slide_helpers(): Initialize helpers for Jupyter notebooks
    - make_figure(): Create consistently-sized matplotlib figures
    - apply_matplotlib_defaults(): Set global chart styling

Basic Usage:
    from deck_builder import (
        SlideContent, DeckBuilder, setup_slide_helpers,
        make_figure, apply_matplotlib_defaults, DECK_CONFIG
    )

    # Configure
    DECK_CONFIG['template_path'] = 'template.pptx'
    DECK_CONFIG['layout_chart'] = [5, 6, 7]

    # Setup
    slides, add_chart_slide, add_section, add_multi_chart_slide = setup_slide_helpers(chart_dir)
    apply_matplotlib_defaults()

    # Create slides
    fig, ax = make_figure('single')
    ax.plot(data)
    add_chart_slide(fig, 'chart.png', 'My Chart')

    # Build deck
    builder = DeckBuilder(DECK_CONFIG['template_path'])
    builder.build(slides, 'output.pptx')

Dependencies:
    pip install python-pptx matplotlib
"""

from collections.abc import Callable
from dataclasses import dataclass
from itertools import cycle
from pathlib import Path

import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# =============================================================================
# GLOBAL CONFIGURATION
# =============================================================================

DECK_CONFIG = {
    # -------------------------------------------------------------------------
    # CHART EXPORT SETTINGS
    # -------------------------------------------------------------------------
    "dpi": 150,
    "fig_facecolor": "white",
    # -------------------------------------------------------------------------
    # FIGURE SIZES (width, height) in inches
    # -------------------------------------------------------------------------
    "fig_single": (12, 6),
    "fig_double": (6, 4),
    "fig_with_kpi": (10, 6),
    "fig_wide": (14, 6),
    "fig_square": (7, 7),
    # -------------------------------------------------------------------------
    # TEMPLATE LAYOUTS - CSI Template12.25.pptx (network: M:\ARS\Presentations\)
    # NOTE: Network template has 17 layouts (0-16). Bundled fallback has fewer.
    # -------------------------------------------------------------------------
    #   0: Cover / Intro - Slide 1      <- Dark bg, centered white text
    #   1: Cover / Intro - Slide 2      <- Title slide (text on left)
    #   2: Divider - Slide 1            <- Section divider
    #   3: Divider - Slide 2
    #   4: Divider - Slide 3            <- Chart with header
    #   5: Analysis - Slide 1           <- Single chart + KPI
    #   6: Analysis - Slide 2           <- Side-by-side (multi_screenshot)
    #   7: Analysis - Slide 3 - Split   <- Side-by-side
    #   8: Analysis - Slide 4           <- Blank, flexible (full width)
    #   9: Analysis - Slide 5           <- Single chart (centered)
    #  10: Analysis - Slide 6 - Main    <- Single chart (right 2/3)
    #  11: Analysis - Slide 7 - Stacked <- Chart on right 2/3
    #  12: Analysis - Slide 8 - Blank   <- Flexible (summary, revenue)
    #  13: Analysis - Slide 9 - Split   <- Full-width spaced (header-only)
    #  14: Agenda                       <- Dedicated Agenda layout
    #  15: Analysis - Slide 11 - Header1
    #  16: Analysis - Slide 11 - Header2
    # -------------------------------------------------------------------------
    "layout_title": 1,
    "layout_title_alt": 0,
    "layout_section": 2,
    "layout_chart": [4, 5, 9, 11],
    "layout_split": [6, 7],
    "layout_stacked": 10,
    "layout_thirds": 8,
    # -------------------------------------------------------------------------
    # LAYOUT-SPECIFIC POSITIONING (in inches)
    # -------------------------------------------------------------------------
    "positioning": {
        "single_default": {
            "left": 0.5,
            "top": 2.5,
            "width": 6,
        },
        "stacked_right": {
            "left": 5.5,
            "top": 2.0,
            "width": 6,
        },
        "split_standard": {
            "top": 2.5,
            "left_pos": 2.3,
            "right_pos": 7.0,
            "width": 4.5,
        },
        "split_spaced": {
            "top": 2.5,
            "left_pos": 0.5,
            "right_pos": 7.5,
            "width": 4.5,
        },
    },
    # -------------------------------------------------------------------------
    # KPI TEXT FORMATTING
    # -------------------------------------------------------------------------
    "kpi_value_size": 28,
    "kpi_label_size": 12,
    # -------------------------------------------------------------------------
    # PATHS
    # -------------------------------------------------------------------------
    "template_path": None,  # Set from ars_config.TEMPLATE_PATH at runtime
}


# =============================================================================
# SLIDE CONTENT DEFINITION
# =============================================================================


@dataclass
class SlideContent:
    """
    Container for all information needed to build a single slide.

    Attributes
    ----------
    slide_type : str
        Options:
            'title'           - Title slide with main title and subtitle
            'section'         - Section divider with large centered text
            'screenshot'      - Single image, nearly full width
            'screenshot_kpi'  - Image on left, KPI callouts on right
            'multi_screenshot'- Two images side by side
            'summary'         - 3x3 grid of bullet points
            'blank'           - Blank placeholder (clears default text)

    title : str
        Text displayed in the slide's title area.

    images : list of str, optional
        File paths to chart/screenshot PNG images.

    kpis : dict, optional
        Key-value pairs for KPI callouts.
        For title slides, use {"subtitle": "January 2025"}

    bullets : list of str, optional
        Text items for summary slides.

    layout_index : int, default 5
        Index of slide layout from template to use.
    """

    slide_type: str
    title: str
    images: list[str] | None = None
    kpis: dict[str, str] | None = None
    bullets: list[str] | None = None
    layout_index: int = 5


# =============================================================================
# DECK BUILDER CLASS
# =============================================================================


class DeckBuilder:
    """
    Assembles SlideContent objects into a PowerPoint presentation.
    """

    # =========================================================================
    # SPACING CONFIGURATION - CSI Template (Template12.25.pptx)
    # Standard PowerPoint slide: 13.33" wide × 7.5" tall
    # =========================================================================

    # Single Screenshot - Default (Layouts 5, 6, 10, 12, 14)
    SINGLE_IMG_LEFT = Inches(0.5)
    SINGLE_IMG_TOP = Inches(2.5)
    SINGLE_IMG_WIDTH = Inches(6)

    # Single Screenshot - Right Side (Layout 11)
    SINGLE_IMG_RIGHT_LEFT = Inches(5.5)
    SINGLE_IMG_RIGHT_TOP = Inches(2.0)
    SINGLE_IMG_RIGHT_WIDTH = Inches(6)

    # Multi-Screenshot - Spaced (Layout 13)
    MULTI_IMG_TOP = Inches(2.5)
    MULTI_IMG_WIDTH = Inches(4.5)
    MULTI_IMG_LEFT_POS = Inches(0.5)
    MULTI_IMG_RIGHT_POS = Inches(7.5)

    # Multi-Screenshot - Standard (Layout 7)
    MULTI_IMG_STD_TOP = Inches(2.5)
    MULTI_IMG_STD_LEFT_POS = Inches(2.3)
    MULTI_IMG_STD_RIGHT_POS = Inches(7.0)
    MULTI_IMG_STD_WIDTH = Inches(4.5)

    # Screenshot with KPIs
    KPI_IMG_LEFT = Inches(0.3)
    KPI_IMG_TOP = Inches(1.5)
    KPI_IMG_WIDTH = Inches(6)

    KPI_TEXT_LEFT = Inches(6.8)
    KPI_TEXT_TOP_START = Inches(1.8)
    KPI_TEXT_WIDTH = Inches(2.5)
    KPI_VALUE_HEIGHT = Inches(0.5)
    KPI_LABEL_HEIGHT = Inches(0.3)
    KPI_SPACING = Inches(1.0)

    # Summary Slide (3x3 bullet grid)
    SUMMARY_COL_POSITIONS = [Inches(0.5), Inches(3.5), Inches(6.5)]
    SUMMARY_ROW_START = Inches(1.8)
    SUMMARY_ROW_SPACING = Inches(1.2)
    SUMMARY_BOX_WIDTH = Inches(2.8)
    SUMMARY_BOX_HEIGHT = Inches(1.0)

    def __init__(self, template_path: str):
        """
        Initialize DeckBuilder with path to PowerPoint template.
        """
        self.template_path = template_path
        self.prs = None

    # Layouts without a title placeholder — need a custom text box
    HEADERONLY_LAYOUTS = {12, 13}

    def _set_title(self, slide, content, title_text, subtitle_text=None):
        """Set slide title, adding a custom text box for header-only layouts."""
        if content.layout_index in self.HEADERONLY_LAYOUTS:
            title_top = Inches(0.30) if content.layout_index == 13 else Inches(0.10)
            tb = slide.shapes.add_textbox(Inches(0.5), title_top, Inches(9.0), Inches(0.6))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title_text
            p.font.size = Pt(24)
            p.font.bold = False
            p.font.color.rgb = RGBColor(255, 255, 255)
            if subtitle_text:
                p2 = tf.add_paragraph()
                p2.text = subtitle_text
                p2.font.size = Pt(12)
                p2.font.color.rgb = RGBColor(220, 220, 220)
        else:
            if slide.shapes.title:
                slide.shapes.title.text = title_text
            if subtitle_text:
                try:
                    slide.placeholders[13].text = subtitle_text
                except (KeyError, IndexError):
                    pass

    def _get_single_positioning(self, layout_index: int) -> tuple:
        """
        Get positioning for single-chart slides based on layout.
        Returns (left, top, width) in Inches.
        Sized for 13.33" wide slides.
        """
        if layout_index == 8:
            return (Inches(0.5), Inches(2.2), Inches(12.0))
        if layout_index == 4:
            return (Inches(2.4), Inches(1.8), Inches(8.5))
        if layout_index == 5:
            return (Inches(2.4), Inches(1.8), Inches(8.5))
        if layout_index == 9:
            return (Inches(2.4), Inches(1.8), Inches(8.5))
        if layout_index == 10:
            return (Inches(5.0), Inches(1.75), Inches(7.8))
        if layout_index == 11:
            return (Inches(2.4), Inches(1.8), Inches(8.5))
        if layout_index == 12:
            return (Inches(0.5), Inches(1.8), Inches(12.0))
        if layout_index == 13:
            return (Inches(0.5), Inches(1.55), Inches(12.0))
        # Default fallback
        return (Inches(2.4), Inches(1.8), Inches(8.5))

    # Maximum chart height: slide is 7.5", title ~1.8", footer ~0.3"
    MAX_CHART_HEIGHT = Inches(5.0)

    def _add_fitted_picture(self, slide, img_path, left, top, max_width, max_height=None):
        """Add image scaled to fit within max_width and max_height."""
        effective_max_h = max_height or self.MAX_CHART_HEIGHT
        from PIL import Image

        with Image.open(img_path) as img:
            native_w, native_h = img.size
        aspect = native_h / native_w
        height_at_width = int(max_width * aspect)
        if height_at_width > effective_max_h:
            slide.shapes.add_picture(img_path, left, top, height=effective_max_h)
        else:
            slide.shapes.add_picture(img_path, left, top, width=max_width)

    def _get_multi_positioning(self, layout_index: int) -> tuple:
        """
        Get positioning for multi-chart slides based on layout.
        Returns (top, left_pos, right_pos, width) in Inches.
        Sized for 13.33" wide slides.
        """
        if layout_index == 6:
            return (Inches(1.8), Inches(0.5), Inches(6.8), Inches(5.8))
        if layout_index == 7:
            return (Inches(1.8), Inches(0.5), Inches(6.8), Inches(5.8))
        # Default
        return (Inches(1.8), Inches(0.5), Inches(6.8), Inches(5.8))

    def build(self, slides: list[SlideContent], output_path: str) -> str:
        """
        Build complete PowerPoint deck from slide definitions.
        """
        self.prs = Presentation(self.template_path)

        for slide_content in slides:
            self._add_slide(slide_content)

        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(output_path)

        return output_path

    def _add_slide(self, content: SlideContent) -> None:
        """
        Create single slide and dispatch to appropriate builder method.
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[content.layout_index])

        builders = {
            "title": self._build_title_slide,
            "section": self._build_section_slide,
            "screenshot": self._build_screenshot_slide,
            "screenshot_kpi": self._build_screenshot_kpi_slide,
            "multi_screenshot": self._build_multi_screenshot_slide,
            "summary": self._build_summary_slide,
            "blank": self._build_blank_slide,
            "mailer_summary": self._build_mailer_summary_slide,
        }

        builder = builders.get(content.slide_type)
        if builder:
            builder(slide, content)
        else:
            raise ValueError(
                f"Unknown slide_type: '{content.slide_type}'. Valid types: {list(builders.keys())}"
            )

    # =========================================================================
    # INDIVIDUAL SLIDE BUILDERS
    # =========================================================================

    def _build_title_slide(self, slide, content: SlideContent) -> None:
        """
        Build title slide with main title and optional subtitle.

        CSI Template cover layouts:
            Layout 0: Cover / Intro - Slide 1 — white text, client name + month
            Layout 1: Cover / Intro - Slide 2 — text box on left
        """
        # -----------------------------------------------------------------
        # Layout 0: Clear all placeholders, add centered white text
        # -----------------------------------------------------------------
        if content.layout_index == 0:
            # Clear ALL placeholder default text
            for ph in slide.placeholders:
                try:
                    for paragraph in ph.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.text = ""
                        paragraph.text = ""
                except Exception:
                    pass

            # Parse title lines (e.g. "Client Name\nJanuary 2026")
            title_lines = content.title.split("\n") if "\n" in content.title else [content.title]

            # Create white text box — centered on dark slide
            text_box = slide.shapes.add_textbox(
                Inches(1.0),  # Left margin
                Inches(2.5),  # Vertically centered area
                Inches(11.0),  # Wide
                Inches(3.0),  # Height
            )

            tf = text_box.text_frame
            tf.word_wrap = True

            for i, line in enumerate(title_lines):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()

                p.text = line
                p.alignment = PP_ALIGN.CENTER

                # Client name = larger, bold
                if i == 0:
                    p.font.size = Pt(38)
                    p.font.bold = True
                else:
                    p.font.size = Pt(26)
                    p.font.bold = False

                # All text white
                p.font.color.rgb = RGBColor(255, 255, 255)

            return

        # -----------------------------------------------------------------
        # Layout 1: Custom text box on left (existing behavior)
        # -----------------------------------------------------------------
        if content.layout_index == 1:
            # Clear ALL placeholder default text so template text doesn't show through
            for ph in slide.placeholders:
                try:
                    for paragraph in ph.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.text = ""
                        paragraph.text = ""
                except Exception:
                    pass

            title_lines = content.title.split("\n") if "\n" in content.title else [content.title]
            subtitle = content.kpis.get("subtitle", "") if content.kpis else ""

            full_text = title_lines[0]
            if len(title_lines) > 1:
                full_text += f"\n{title_lines[1]}"
            if len(title_lines) > 2:
                full_text += f"\n{title_lines[2]}"
            elif subtitle:
                full_text += f"\n{'─' * 20}\n{subtitle}"

            text_box = slide.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(6.0), Inches(2.0))

            tf = text_box.text_frame
            tf.word_wrap = True

            lines = full_text.split("\n")
            for i, line in enumerate(lines):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()

                p.text = line
                p.alignment = PP_ALIGN.LEFT
                p.font.color.rgb = RGBColor(255, 255, 255)

                if i == 0:
                    p.font.size = Pt(34)
                    p.font.bold = True
                else:
                    p.font.size = Pt(20)

            return

        # -----------------------------------------------------------------
        # Other layouts: Use placeholders
        # -----------------------------------------------------------------
        title_lines = content.title.split("\n") if "\n" in content.title else [content.title]

        if slide.shapes.title:
            slide.shapes.title.text = title_lines[0]

        additional_text = title_lines[1:] if len(title_lines) > 1 else []

        if content.kpis and "subtitle" in content.kpis:
            additional_text.append(content.kpis["subtitle"])

        text_placeholders = [26, 29, 30, 27, 28, 31]

        for i, text in enumerate(additional_text):
            if i < len(text_placeholders):
                try:
                    ph = slide.placeholders[text_placeholders[i]]
                    ph.text = text
                    break
                except (KeyError, IndexError):
                    try:
                        slide.placeholders[1].text = text
                    except (KeyError, IndexError):
                        pass

    def _build_section_slide(self, slide, content: SlideContent) -> None:
        """
        Build section divider slide.

        CSI Template divider layouts:
            - Layout 2: Title at top, multiple text areas below
            - Layout 3: Title on right
            - Layout 4: Title at top
        """
        title_text = content.title
        subtitle_text = None

        if "\n" in content.title:
            parts = content.title.split("\n", 1)
            title_text = parts[0]
            subtitle_text = parts[1] if len(parts) > 1 else None

        if slide.shapes.title:
            slide.shapes.title.text = title_text

        if subtitle_text:
            for ph_idx in [44, 13, 1, 14]:
                try:
                    ph = slide.placeholders[ph_idx]
                    ph.text = subtitle_text
                    break
                except (KeyError, IndexError):
                    continue

    def _build_blank_slide(self, slide, content: SlideContent) -> None:
        """
        Build blank placeholder slide.

        Creates a slide using the specified layout, clears all default
        placeholder text, then optionally sets the title. Used for slides
        where James adds content manually (executive summary, revenue, etc).
        """
        # Clear ALL placeholder default text
        for ph in slide.placeholders:
            try:
                for paragraph in ph.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = ""
                    paragraph.text = ""
            except Exception:
                pass

        # Layout 8: handle subtitle (text after \n) separately
        # Keep the subtitle placeholder (gold line area) if we have subtitle text
        title_text = content.title or ""
        subtitle_text = None
        if "\n" in title_text:
            parts = title_text.split("\n", 1)
            title_text = parts[0]
            subtitle_text = parts[1]

        if content.layout_index == 8:
            if subtitle_text:
                # Keep subtitle placeholder (idx 1) for the gold-line subtitle area
                to_remove = [
                    ph for ph in slide.placeholders
                    if ph.placeholder_format.idx not in (0, 1)
                ]
            else:
                # No subtitle — remove all non-title placeholders
                to_remove = [
                    ph for ph in slide.placeholders
                    if ph.placeholder_format.idx != 0
                ]
            for ph in to_remove:
                ph.element.getparent().remove(ph.element)

        # Set the title
        if title_text:
            if content.layout_index in self.HEADERONLY_LAYOUTS:
                self._set_title(slide, content, title_text, subtitle_text)
            elif slide.shapes.title:
                slide.shapes.title.text = title_text
                # Layout 0 has dark background — make title visible
                if content.layout_index == 0:
                    for p in slide.shapes.title.text_frame.paragraphs:
                        p.font.color.rgb = RGBColor(255, 255, 255)
                        p.font.size = Pt(28)
                        p.font.bold = True
                        p.alignment = PP_ALIGN.LEFT

        # Set subtitle in placeholder idx 1 (next to gold line on layout 8)
        if subtitle_text and content.layout_index == 8:
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 1:
                    ph.text = subtitle_text
                    break

    def _build_mailer_summary_slide(self, slide, content: SlideContent) -> None:
        """
        Build composite mailer summary slide — 13.33" wide × 7.5" tall.

        Layout (3 equal columns):
          Row 1: Title (from template header)
          Row 2: Insight text (left) | Mailer KPIs (right, 3 across)
          Row 3: Donut chart  |  Horizontal bar chart  |  Inside the Numbers

        All section headers use unified Pt(16) bold.
        """
        # Clear non-title placeholders to prevent default text showing through
        for ph in slide.placeholders:
            if ph.placeholder_format.idx != 0:
                try:
                    ph.element.getparent().remove(ph.element)
                except Exception:
                    pass

        # =====================================================================
        # COLUMN GEOMETRY — 3 equal columns on 13.33" slide
        # =====================================================================
        COL1_L = Inches(0.2)  # donut
        COL2_L = Inches(4.4)  # bar chart
        COL3_L = Inches(8.8)  # inside the numbers
        COL_W = Inches(4.1)  # each column width
        HEADER_SIZE = Pt(16)  # unified section header font

        # Vertical geometry — insight + KPIs sit below title, then charts
        ROW1_TOP = Inches(1.6)  # insight text + KPI header
        KPI_VAL_TOP = Inches(2.1)  # KPI values
        KPI_LBL_TOP = Inches(2.55)  # KPI labels
        SECT_TOP = Inches(3.2)  # section header row
        CHART_TOP = Inches(3.5)  # top of chart images

        # --- TITLE (custom position for mailer summary, no subtitle) ---
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.38), Inches(9.0), Inches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content.title
        p.font.size = Pt(24)
        p.font.bold = False
        p.font.color.rgb = RGBColor(255, 255, 255)

        # --- Parse bullets ---
        insight_text = ""
        inside_numbers = []
        if content.bullets:
            insight_text = content.bullets[0] if content.bullets[0] else ""
            inside_numbers = content.bullets[1:]

        # --- INSIGHT TEXT (upper-left, top-aligned with KPIs) ---
        if insight_text:
            tb = slide.shapes.add_textbox(Inches(0.5), ROW1_TOP, Inches(5.0), Inches(1.0))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = insight_text
            p.font.size = Pt(15)
            p.font.color.rgb = RGBColor(0, 0, 0)

        # --- MAILER KPIs (upper-right, top-aligned with insight) ---
        if content.kpis:
            kpi_left = Inches(7.8)
            kpi_block_w = Inches(5.2)
            tb = slide.shapes.add_textbox(kpi_left, ROW1_TOP, kpi_block_w, Inches(0.4))
            tf = tb.text_frame
            p = tf.paragraphs[0]
            p.text = "Mailer KPIs"
            p.font.size = HEADER_SIZE
            p.font.bold = True
            p.font.color.rgb = RGBColor(30, 61, 89)
            p.alignment = PP_ALIGN.CENTER

            kpi_items = list(content.kpis.items())
            kpi_each_w = 5.2 / max(len(kpi_items), 1)

            for i, (label, value) in enumerate(kpi_items):
                x = 7.8 + i * kpi_each_w

                tb = slide.shapes.add_textbox(
                    Inches(x), KPI_VAL_TOP, Inches(kpi_each_w), Inches(0.5)
                )
                tf = tb.text_frame
                p = tf.paragraphs[0]
                p.text = str(value)
                p.font.size = Pt(24)
                p.font.bold = True
                p.font.color.rgb = RGBColor(30, 61, 89)
                p.alignment = PP_ALIGN.CENTER

                tb = slide.shapes.add_textbox(
                    Inches(x), KPI_LBL_TOP, Inches(kpi_each_w), Inches(0.3)
                )
                tf = tb.text_frame
                p = tf.paragraphs[0]
                p.text = label
                p.font.size = Pt(13)
                p.font.color.rgb = RGBColor(0, 0, 0)
                p.alignment = PP_ALIGN.CENTER

        # --- SECTION HEADERS (3 aligned across columns) ---
        for col_left, header_text in [
            (COL1_L, "Response Share"),
            (COL2_L, "Response Rate"),
            (COL3_L, "Inside the Numbers"),
        ]:
            tb = slide.shapes.add_textbox(col_left, SECT_TOP, COL_W, Inches(0.3))
            tf = tb.text_frame
            p = tf.paragraphs[0]
            p.text = header_text
            p.font.size = HEADER_SIZE
            p.font.bold = True
            p.font.color.rgb = RGBColor(30, 61, 89)
            p.alignment = PP_ALIGN.CENTER

        # --- DONUT CHART (column 1) ---
        if content.images and len(content.images) > 0 and Path(content.images[0]).exists():
            slide.shapes.add_picture(content.images[0], COL1_L, CHART_TOP, width=COL_W)

        # --- HORIZONTAL BAR CHART (column 2) ---
        if content.images and len(content.images) > 1 and Path(content.images[1]).exists():
            slide.shapes.add_picture(content.images[1], COL2_L, CHART_TOP, width=COL_W)

        # --- INSIDE THE NUMBERS (column 3, below header) ---
        if inside_numbers:
            for i, item in enumerate(inside_numbers):
                if "|" in item:
                    pct, desc = item.split("|", 1)
                else:
                    pct, desc = item, ""

                y_pos = 3.9 + i * 1.2

                # Percentage + description side by side, vertically centered
                tb = slide.shapes.add_textbox(COL3_L, Inches(y_pos), Inches(1.4), Inches(0.9))
                tf = tb.text_frame
                tf.auto_size = None
                tf.word_wrap = False
                try:
                    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                    bodyPr = tf._txBody.find(qn("a:bodyPr"))
                    bodyPr.set("anchor", "ctr")
                except Exception:
                    pass
                p = tf.paragraphs[0]
                p.text = pct.strip()
                p.font.size = Pt(26)
                p.font.bold = True
                p.font.color.rgb = RGBColor(0, 80, 114)

                if desc:
                    tb = slide.shapes.add_textbox(
                        Inches(10.2), Inches(y_pos), Inches(2.6), Inches(0.9)
                    )
                    tf = tb.text_frame
                    tf.auto_size = None
                    tf.word_wrap = True
                    try:
                        bodyPr = tf._txBody.find(qn("a:bodyPr"))
                        bodyPr.set("anchor", "ctr")
                    except Exception:
                        pass
                    p = tf.paragraphs[0]
                    p.text = desc.strip()
                    p.font.size = Pt(13)
                    p.font.color.rgb = RGBColor(0, 0, 0)

    def _build_screenshot_slide(self, slide, content: SlideContent) -> None:
        """
        Build slide with title, subtitle, and single image.
        """
        title_text = content.title
        subtitle_text = None

        if "\n" in content.title:
            parts = content.title.split("\n", 1)
            title_text = parts[0]
            subtitle_text = parts[1] if len(parts) > 1 else None

        # Remove non-title placeholders on layout 8 to prevent gold bar
        if content.layout_index == 8:
            to_remove = [ph for ph in slide.placeholders if ph.placeholder_format.idx != 0]
            for ph in to_remove:
                ph.element.getparent().remove(ph.element)

        self._set_title(slide, content, title_text, subtitle_text)

        left, top, width = self._get_single_positioning(content.layout_index)
        extra_h = Inches(5.2) if content.layout_index in (8, 12, 13) else None

        if content.images and Path(content.images[0]).exists():
            self._add_fitted_picture(slide, content.images[0], left, top, width, max_height=extra_h)

    def _build_screenshot_kpi_slide(self, slide, content: SlideContent) -> None:
        """
        Build slide with image and KPI callouts.
        """
        title_text = content.title
        subtitle_text = None

        if "\n" in content.title:
            parts = content.title.split("\n", 1)
            title_text = parts[0]
            subtitle_text = parts[1] if len(parts) > 1 else None

        # Remove non-title placeholders on layout 8 to prevent gold bar
        if content.layout_index == 8:
            to_remove = [ph for ph in slide.placeholders if ph.placeholder_format.idx != 0]
            for ph in to_remove:
                ph.element.getparent().remove(ph.element)

        self._set_title(slide, content, title_text, subtitle_text)

        left, top, width = self._get_single_positioning(content.layout_index)
        extra_h = Inches(5.2) if content.layout_index in (8, 12, 13) else None

        if content.images and Path(content.images[0]).exists():
            self._add_fitted_picture(slide, content.images[0], left, top, width, max_height=extra_h)

        if content.kpis:
            kpi_placeholder_pairs = [
                (26, 19),  # First KPI
                (27, 28),  # Second KPI
            ]

            kpi_items = [(k, v) for k, v in content.kpis.items() if k != "subtitle"]

            for i, (label, value) in enumerate(kpi_items):
                if i >= len(kpi_placeholder_pairs):
                    break

                label_idx, value_idx = kpi_placeholder_pairs[i]

                try:
                    slide.placeholders[label_idx].text = label
                except (KeyError, IndexError):
                    pass

                try:
                    slide.placeholders[value_idx].text = str(value)
                except (KeyError, IndexError):
                    pass

    def _build_multi_screenshot_slide(self, slide, content: SlideContent) -> None:
        """
        Build slide with two images side by side.
        """
        title_text = content.title
        subtitle_text = None

        if "\n" in content.title:
            parts = content.title.split("\n", 1)
            title_text = parts[0]
            subtitle_text = parts[1] if len(parts) > 1 else None

        self._set_title(slide, content, title_text, subtitle_text)

        top, left_pos, right_pos, width = self._get_multi_positioning(content.layout_index)

        positions = [
            (left_pos, top, width),
            (right_pos, top, width),
        ]

        if content.images:
            for i, img_path in enumerate(content.images[:2]):
                if Path(img_path).exists():
                    left, img_top, img_width = positions[i]
                    self._add_fitted_picture(
                        slide, img_path, left, img_top, img_width, max_height=Inches(5.5)
                    )

    def _build_summary_slide(self, slide, content: SlideContent) -> None:
        """
        Build summary slide with bullets in 3x3 grid.
        """
        if slide.shapes.title:
            slide.shapes.title.text = content.title

        if content.bullets:
            for i, bullet_text in enumerate(content.bullets[:9]):
                col = i % 3
                row = i // 3

                left = self.SUMMARY_COL_POSITIONS[col]
                top = self.SUMMARY_ROW_START + (row * self.SUMMARY_ROW_SPACING)

                text_box = slide.shapes.add_textbox(
                    left, top, self.SUMMARY_BOX_WIDTH, self.SUMMARY_BOX_HEIGHT
                )

                tf = text_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = bullet_text
                p.font.size = Pt(11)


# =============================================================================
# UTILITY FUNCTIONS
# =============================================================================


def inspect_template(template_path: str) -> None:
    """
    Print all available slide layouts in a PowerPoint template.
    """
    prs = Presentation(template_path)
    print(f"Layouts in {template_path}:\n")
    for i, layout in enumerate(prs.slide_layouts):
        print(f"  {i}: {layout.name}")


def create_title_slide(client_name: str, period: str, layout_index: int = None) -> SlideContent:
    """
    Convenience function to create a title slide.
    """
    actual_layout = layout_index if layout_index is not None else DECK_CONFIG.get("layout_title", 1)
    return SlideContent(
        slide_type="title",
        title=f"{client_name} Monthly Report",
        kpis={"subtitle": period},
        layout_index=actual_layout,
    )


def create_summary_slide(
    title: str,
    col1_bullets: list[str],
    col2_bullets: list[str],
    col3_bullets: list[str],
    layout_index: int = None,
) -> SlideContent:
    """
    Convenience function to create 3-column summary slide.
    """
    bullets = []
    for i in range(3):
        bullets.append(col1_bullets[i] if i < len(col1_bullets) else "")
        bullets.append(col2_bullets[i] if i < len(col2_bullets) else "")
        bullets.append(col3_bullets[i] if i < len(col3_bullets) else "")

    actual_layout = layout_index
    if actual_layout is None:
        chart_layouts = DECK_CONFIG.get("layout_chart", [5])
        actual_layout = chart_layouts[0] if isinstance(chart_layouts, list) else chart_layouts

    return SlideContent(
        slide_type="summary", title=title, bullets=bullets, layout_index=actual_layout
    )


# =============================================================================
# FIGURE CREATION
# =============================================================================


def make_figure(fig_type: str = "single") -> tuple:
    """
    Create a matplotlib figure with standardized sizing.
    """
    size_map = {
        "single": "fig_single",
        "double": "fig_double",
        "kpi": "fig_with_kpi",
        "wide": "fig_wide",
        "square": "fig_square",
    }

    config_key = size_map.get(fig_type, "fig_single")
    figsize = DECK_CONFIG.get(config_key, (10, 6))

    return plt.subplots(figsize=figsize)


# =============================================================================
# MATPLOTLIB DEFAULTS
# =============================================================================


def apply_matplotlib_defaults() -> None:
    """
    Apply consistent matplotlib styling to all charts.
    """
    plt.rcParams.update(
        {
            "font.size": 12,
            "axes.titlesize": 14,
            "axes.labelsize": 12,
            "xtick.labelsize": 10,
            "ytick.labelsize": 10,
            "legend.fontsize": 10,
            "figure.titlesize": 16,
            "axes.spines.top": False,
            "axes.spines.right": False,
            "axes.grid": False,
            "figure.facecolor": "white",
            "axes.facecolor": "white",
            "savefig.facecolor": "white",
            "savefig.bbox": "tight",
        }
    )


# =============================================================================
# JUPYTER NOTEBOOK HELPER SETUP
# =============================================================================


def setup_slide_helpers(
    chart_dir: Path, deck_config: dict = None
) -> tuple[list, Callable, Callable, Callable]:
    """
    Initialize slide list and helper functions for Jupyter notebook.

    Returns
    -------
    tuple of (slides, add_chart_slide, add_section, add_multi_chart_slide)
    """
    config = deck_config or DECK_CONFIG

    chart_dir = Path(chart_dir)
    chart_dir.mkdir(parents=True, exist_ok=True)

    slides = []

    chart_layouts = config.get("layout_chart", [5])
    if not isinstance(chart_layouts, list):
        chart_layouts = [chart_layouts]
    layout_cycler = cycle(chart_layouts)

    dpi = config.get("dpi", 150)
    facecolor = config.get("fig_facecolor", "white")

    # -------------------------------------------------------------------------
    # Helper: add_chart_slide
    # -------------------------------------------------------------------------
    def add_chart_slide(
        fig,
        filename: str,
        title: str,
        slide_type: str = "screenshot",
        kpis: dict = None,
        layout_index: int = None,
    ) -> None:
        """Save matplotlib figure and add slide to list."""
        path = chart_dir / filename
        fig.savefig(path, dpi=dpi, bbox_inches="tight", facecolor=facecolor)
        plt.close(fig)

        actual_layout = layout_index if layout_index is not None else next(layout_cycler)

        slides.append(
            SlideContent(
                slide_type=slide_type,
                title=title,
                images=[str(path)],
                kpis=kpis,
                layout_index=actual_layout,
            )
        )

    # -------------------------------------------------------------------------
    # Helper: add_section
    # -------------------------------------------------------------------------
    def add_section(title: str, layout_index: int = None) -> None:
        """Add section divider slide."""
        actual_layout = (
            layout_index if layout_index is not None else config.get("layout_section", 2)
        )
        slides.append(SlideContent("section", title, layout_index=actual_layout))

    # -------------------------------------------------------------------------
    # Helper: add_multi_chart_slide
    # -------------------------------------------------------------------------
    def add_multi_chart_slide(
        fig1, fig2, filename1: str, filename2: str, title: str, layout_index: int = None
    ) -> None:
        """Save two figures and add as side-by-side slide."""
        path1 = chart_dir / filename1
        path2 = chart_dir / filename2

        fig1.savefig(path1, dpi=dpi, bbox_inches="tight", facecolor=facecolor)
        fig2.savefig(path2, dpi=dpi, bbox_inches="tight", facecolor=facecolor)

        plt.close(fig1)
        plt.close(fig2)

        actual_layout = layout_index if layout_index is not None else next(layout_cycler)

        slides.append(
            SlideContent(
                slide_type="multi_screenshot",
                title=title,
                images=[str(path1), str(path2)],
                layout_index=actual_layout,
            )
        )

    return slides, add_chart_slide, add_section, add_multi_chart_slide


# =============================================================================
# MODULE TEST / EXAMPLE
# =============================================================================

if __name__ == "__main__":
    print("deck_builder.py - PowerPoint Automation Module")
    print("=" * 50)
    print()
    print("Usage:")
    print("  from deck_builder import (")
    print("      SlideContent, DeckBuilder, setup_slide_helpers,")
    print("      make_figure, apply_matplotlib_defaults, DECK_CONFIG")
    print("  )")
    print()
    print("Quick Start:")
    print("  1. Run: inspect_template('your_template.pptx')")
    print("  2. Update DECK_CONFIG['layout_chart'] with valid indices")
    print("  3. Call setup_slide_helpers(chart_dir)")
    print("  4. Use add_chart_slide() after each analysis")
    print("  5. Call builder.build(slides, 'output.pptx')")
    print()
    print("See README.md for full documentation.")
