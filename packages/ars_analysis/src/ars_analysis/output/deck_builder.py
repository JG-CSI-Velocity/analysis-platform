"""PowerPoint deck builder -- ported from original ars_analysis-jupyter/deck_builder.py.

This module restores the full production-quality PPTX generation with:
- 8 slide types (title, section, screenshot, screenshot_kpi, multi_screenshot, etc.)
- Layout-specific positioning for 10+ template layouts
- Aspect-ratio preserving image scaling via PIL
- Preamble slides (13 intro/section/placeholder slides)
- Consolidation logic (merge paired slides, separate appendix)
- Section ordering matching the original ARS reference deck
"""

from __future__ import annotations

import calendar
from dataclasses import dataclass
from pathlib import Path

from loguru import logger
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from ars_analysis.pipeline.context import PipelineContext

# Embedded fallback template (ships with the package)
_FALLBACK_TEMPLATE = Path(__file__).parent / "template" / "Template12.25.pptx"


# =============================================================================
# SLIDE CONTENT DEFINITION
# =============================================================================


@dataclass
class SlideContent:
    """Container for all information needed to build a single slide.

    slide_type options:
        'title'            - Title slide with main title and subtitle
        'section'          - Section divider with large centered text
        'screenshot'       - Single image, nearly full width
        'screenshot_kpi'   - Image on left, KPI callouts on right
        'multi_screenshot' - Two images side by side
        'summary'          - 3x3 grid of bullet points
        'blank'            - Blank placeholder (clears default text)
        'mailer_summary'   - 3-column composite (donut, bar, inside numbers)
    """

    slide_type: str
    title: str
    images: list[str] | None = None
    kpis: dict[str, str] | None = None
    bullets: list[str] | None = None
    layout_index: int = 5


# =============================================================================
# DECK BUILDER CLASS -- ported from original ars_analysis-jupyter/deck_builder.py
# =============================================================================


class DeckBuilder:
    """Assembles SlideContent objects into a PowerPoint presentation."""

    # Spacing for CSI Template (Template12.25.pptx)
    # Standard PowerPoint slide: 13.33" wide x 7.5" tall

    # Single Screenshot - Default (Layouts 5, 6, 10, 12, 14)
    SINGLE_IMG_LEFT = Inches(0.5)
    SINGLE_IMG_TOP = Inches(2.5)
    SINGLE_IMG_WIDTH = Inches(6)

    # Single Screenshot - Right Side (Layout 11)
    SINGLE_IMG_RIGHT_LEFT = Inches(5.5)
    SINGLE_IMG_RIGHT_TOP = Inches(2.0)
    SINGLE_IMG_RIGHT_WIDTH = Inches(6)

    # Multi-Screenshot - Spaced (Layout 13)
    MULTI_IMG_TOP = Inches(2.5)
    MULTI_IMG_WIDTH = Inches(4.5)
    MULTI_IMG_LEFT_POS = Inches(0.5)
    MULTI_IMG_RIGHT_POS = Inches(7.5)

    # Multi-Screenshot - Standard (Layout 7)
    MULTI_IMG_STD_TOP = Inches(2.5)
    MULTI_IMG_STD_LEFT_POS = Inches(2.3)
    MULTI_IMG_STD_RIGHT_POS = Inches(7.0)
    MULTI_IMG_STD_WIDTH = Inches(4.5)

    # Screenshot with KPIs
    KPI_IMG_LEFT = Inches(0.3)
    KPI_IMG_TOP = Inches(1.5)
    KPI_IMG_WIDTH = Inches(6)

    KPI_TEXT_LEFT = Inches(6.8)
    KPI_TEXT_TOP_START = Inches(1.8)
    KPI_TEXT_WIDTH = Inches(2.5)
    KPI_VALUE_HEIGHT = Inches(0.5)
    KPI_LABEL_HEIGHT = Inches(0.3)
    KPI_SPACING = Inches(1.0)

    # Summary Slide (3x3 bullet grid)
    SUMMARY_COL_POSITIONS = [Inches(0.5), Inches(3.5), Inches(6.5)]
    SUMMARY_ROW_START = Inches(1.8)
    SUMMARY_ROW_SPACING = Inches(1.2)
    SUMMARY_BOX_WIDTH = Inches(2.8)
    SUMMARY_BOX_HEIGHT = Inches(1.0)

    # Maximum chart height
    MAX_CHART_HEIGHT = Inches(5.0)

    # Layouts without a title placeholder -- need a custom text box
    HEADERONLY_LAYOUTS = {12, 13}

    def __init__(self, template_path: str):
        self.template_path = template_path
        self.prs = None

    def build(self, slides: list[SlideContent], output_path: str) -> str:
        """Build complete PowerPoint deck from slide definitions."""
        self.prs = Presentation(self.template_path)
        n_layouts = len(self.prs.slide_layouts)

        for i, slide_content in enumerate(slides):
            if slide_content.layout_index >= n_layouts:
                logger.warning(
                    "Slide {i} '{title}' has layout_index={idx} but template only has {n} layouts, using 0",
                    i=i,
                    title=slide_content.title[:40],
                    idx=slide_content.layout_index,
                    n=n_layouts,
                )
                slide_content.layout_index = 0
            try:
                self._add_slide(slide_content)
            except Exception as exc:
                logger.error(
                    "Slide {i} '{title}' (type={t}, layout={l}) failed: {err}",
                    i=i,
                    title=slide_content.title[:40],
                    t=slide_content.slide_type,
                    l=slide_content.layout_index,
                    err=exc,
                )

        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(output_path)
        return output_path

    def _add_slide(self, content: SlideContent) -> None:
        """Create single slide and dispatch to appropriate builder method."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[content.layout_index])

        builders = {
            "title": self._build_title_slide,
            "section": self._build_section_slide,
            "screenshot": self._build_screenshot_slide,
            "screenshot_kpi": self._build_screenshot_kpi_slide,
            "multi_screenshot": self._build_multi_screenshot_slide,
            "summary": self._build_summary_slide,
            "blank": self._build_blank_slide,
            "mailer_summary": self._build_mailer_summary_slide,
        }

        builder = builders.get(content.slide_type)
        if builder:
            builder(slide, content)
        else:
            logger.warning("Unknown slide_type: {t}", t=content.slide_type)

    # -------------------------------------------------------------------------
    # Title helpers
    # -------------------------------------------------------------------------

    def _set_title(self, slide, content, title_text, subtitle_text=None):
        """Set slide title, adding custom text box for header-only layouts."""
        if content.layout_index in self.HEADERONLY_LAYOUTS:
            title_top = Inches(0.30) if content.layout_index == 13 else Inches(0.10)
            tb = slide.shapes.add_textbox(Inches(0.5), title_top, Inches(9.0), Inches(0.6))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title_text
            p.font.size = Pt(24)
            p.font.bold = False
            p.font.color.rgb = RGBColor(255, 255, 255)
            if subtitle_text:
                p2 = tf.add_paragraph()
                p2.text = subtitle_text
                p2.font.size = Pt(12)
                p2.font.color.rgb = RGBColor(220, 220, 220)
        else:
            if slide.shapes.title:
                slide.shapes.title.text = title_text
            if subtitle_text:
                try:
                    slide.placeholders[13].text = subtitle_text
                except (KeyError, IndexError):
                    pass

    def _get_single_positioning(self, layout_index: int) -> tuple:
        """Get (left, top, width) positioning for single-chart slides."""
        if layout_index == 8:
            return (Inches(0.5), Inches(2.2), Inches(12.0))
        if layout_index in (4, 5, 9, 11):
            return (Inches(2.4), Inches(1.8), Inches(8.5))
        if layout_index == 10:
            return (Inches(5.0), Inches(1.75), Inches(7.8))
        if layout_index == 12:
            return (Inches(0.5), Inches(1.8), Inches(12.0))
        if layout_index == 13:
            return (Inches(0.5), Inches(1.55), Inches(12.0))
        return (Inches(2.4), Inches(1.8), Inches(8.5))

    def _get_multi_positioning(self, layout_index: int) -> tuple:
        """Get (top, left_pos, right_pos, width) for multi-chart slides."""
        if layout_index in (6, 7):
            return (Inches(1.2), Inches(0.5), Inches(6.8), Inches(5.8))
        return (Inches(1.2), Inches(0.5), Inches(6.8), Inches(5.8))

    def _add_fitted_picture(self, slide, img_path, left, top, max_width, max_height=None):
        """Add image scaled to fit within max_width and max_height."""
        effective_max_h = max_height or self.MAX_CHART_HEIGHT
        try:
            from PIL import Image

            with Image.open(img_path) as img:
                native_w, native_h = img.size
            aspect = native_h / native_w
            height_at_width = int(max_width * aspect)
            if height_at_width > effective_max_h:
                slide.shapes.add_picture(img_path, left, top, height=effective_max_h)
            else:
                slide.shapes.add_picture(img_path, left, top, width=max_width)
        except ImportError:
            slide.shapes.add_picture(img_path, left, top, width=max_width)

    # -------------------------------------------------------------------------
    # Individual slide builders
    # -------------------------------------------------------------------------

    def _build_title_slide(self, slide, content: SlideContent) -> None:
        """Build title slide with main title and optional subtitle."""
        if content.layout_index == 0:
            for ph in slide.placeholders:
                try:
                    for paragraph in ph.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.text = ""
                        paragraph.text = ""
                except Exception:
                    pass

            title_lines = content.title.split("\n") if "\n" in content.title else [content.title]
            text_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.0), Inches(3.0))
            tf = text_box.text_frame
            tf.word_wrap = True

            for i, line in enumerate(title_lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = line
                p.alignment = PP_ALIGN.CENTER
                if i == 0:
                    p.font.size = Pt(38)
                    p.font.bold = True
                else:
                    p.font.size = Pt(26)
                    p.font.bold = False
                p.font.color.rgb = RGBColor(255, 255, 255)
            return

        if content.layout_index == 1:
            title_lines = content.title.split("\n") if "\n" in content.title else [content.title]
            subtitle = content.kpis.get("subtitle", "") if content.kpis else ""

            full_text = title_lines[0]
            if len(title_lines) > 1:
                full_text += f"\n{title_lines[1]}"
            if len(title_lines) > 2:
                full_text += f"\n{title_lines[2]}"
            elif subtitle:
                full_text += f"\n{'─' * 20}\n{subtitle}"

            text_box = slide.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(6.0), Inches(2.0))
            tf = text_box.text_frame
            tf.word_wrap = True

            lines = full_text.split("\n")
            for i, line in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = line
                p.alignment = PP_ALIGN.LEFT
                p.font.color.rgb = RGBColor(255, 255, 255)
                if i == 0:
                    p.font.size = Pt(34)
                    p.font.bold = True
                else:
                    p.font.size = Pt(20)
            return

        title_lines = content.title.split("\n") if "\n" in content.title else [content.title]
        if slide.shapes.title:
            slide.shapes.title.text = title_lines[0]

        additional_text = title_lines[1:] if len(title_lines) > 1 else []
        if content.kpis and "subtitle" in content.kpis:
            additional_text.append(content.kpis["subtitle"])

        text_placeholders = [26, 29, 30, 27, 28, 31]
        for i, text in enumerate(additional_text):
            if i < len(text_placeholders):
                try:
                    slide.placeholders[text_placeholders[i]].text = text
                    break
                except (KeyError, IndexError):
                    try:
                        slide.placeholders[1].text = text
                    except (KeyError, IndexError):
                        pass

    def _build_section_slide(self, slide, content: SlideContent) -> None:
        """Build section divider slide."""
        title_text = content.title
        subtitle_text = None

        if "\n" in content.title:
            parts = content.title.split("\n", 1)
            title_text = parts[0]
            subtitle_text = parts[1] if len(parts) > 1 else None

        if slide.shapes.title:
            slide.shapes.title.text = title_text

        if subtitle_text:
            for ph_idx in [44, 13, 1, 14]:
                try:
                    slide.placeholders[ph_idx].text = subtitle_text
                    break
                except (KeyError, IndexError):
                    continue

    def _build_blank_slide(self, slide, content: SlideContent) -> None:
        """Build blank placeholder slide."""
        for ph in slide.placeholders:
            try:
                for paragraph in ph.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = ""
                    paragraph.text = ""
            except Exception:
                pass

        if content.layout_index == 8:
            to_remove = [ph for ph in slide.placeholders if ph.placeholder_format.idx != 0]
            for ph in to_remove:
                ph.element.getparent().remove(ph.element)

        if content.title:
            if content.layout_index in self.HEADERONLY_LAYOUTS:
                self._set_title(slide, content, content.title)
            elif slide.shapes.title:
                slide.shapes.title.text = content.title
                if content.layout_index == 0:
                    for p in slide.shapes.title.text_frame.paragraphs:
                        p.font.color.rgb = RGBColor(255, 255, 255)
                        p.font.size = Pt(28)
                        p.font.bold = True
                        p.alignment = PP_ALIGN.LEFT

    def _build_screenshot_slide(self, slide, content: SlideContent) -> None:
        """Build slide with title, subtitle, and single image."""
        title_text = content.title
        subtitle_text = None

        if "\n" in content.title:
            parts = content.title.split("\n", 1)
            title_text = parts[0]
            subtitle_text = parts[1] if len(parts) > 1 else None

        if content.layout_index == 8:
            to_remove = [ph for ph in slide.placeholders if ph.placeholder_format.idx != 0]
            for ph in to_remove:
                ph.element.getparent().remove(ph.element)

        self._set_title(slide, content, title_text, subtitle_text)

        left, top, width = self._get_single_positioning(content.layout_index)
        extra_h = Inches(5.2) if content.layout_index in (8, 12, 13) else None

        if content.images and Path(content.images[0]).exists():
            self._add_fitted_picture(slide, content.images[0], left, top, width, max_height=extra_h)

    def _build_screenshot_kpi_slide(self, slide, content: SlideContent) -> None:
        """Build slide with image and KPI callouts."""
        title_text = content.title
        subtitle_text = None

        if "\n" in content.title:
            parts = content.title.split("\n", 1)
            title_text = parts[0]
            subtitle_text = parts[1] if len(parts) > 1 else None

        if content.layout_index == 8:
            to_remove = [ph for ph in slide.placeholders if ph.placeholder_format.idx != 0]
            for ph in to_remove:
                ph.element.getparent().remove(ph.element)

        self._set_title(slide, content, title_text, subtitle_text)

        left, top, width = self._get_single_positioning(content.layout_index)
        extra_h = Inches(5.2) if content.layout_index in (8, 12, 13) else None

        if content.images and Path(content.images[0]).exists():
            self._add_fitted_picture(slide, content.images[0], left, top, width, max_height=extra_h)

        if content.kpis:
            kpi_placeholder_pairs = [
                (26, 19),
                (27, 28),
            ]
            kpi_items = [(k, v) for k, v in content.kpis.items() if k != "subtitle"]
            for i, (label_text, value) in enumerate(kpi_items):
                if i >= len(kpi_placeholder_pairs):
                    break
                label_idx, value_idx = kpi_placeholder_pairs[i]
                try:
                    slide.placeholders[label_idx].text = label_text
                except (KeyError, IndexError):
                    pass
                try:
                    slide.placeholders[value_idx].text = str(value)
                except (KeyError, IndexError):
                    pass

    def _build_multi_screenshot_slide(self, slide, content: SlideContent) -> None:
        """Build slide with two images side by side."""
        title_text = content.title
        subtitle_text = None

        if "\n" in content.title:
            parts = content.title.split("\n", 1)
            title_text = parts[0]
            subtitle_text = parts[1] if len(parts) > 1 else None

        self._set_title(slide, content, title_text, subtitle_text)

        top, left_pos, right_pos, width = self._get_multi_positioning(content.layout_index)

        positions = [
            (left_pos, top, width),
            (right_pos, top, width),
        ]

        if content.images:
            for i, img_path in enumerate(content.images[:2]):
                if Path(img_path).exists():
                    left, img_top, img_width = positions[i]
                    self._add_fitted_picture(
                        slide, img_path, left, img_top, img_width, max_height=Inches(5.5)
                    )

    def _build_summary_slide(self, slide, content: SlideContent) -> None:
        """Build summary slide with bullets in 3x3 grid."""
        if slide.shapes.title:
            slide.shapes.title.text = content.title

        if content.bullets:
            for i, bullet_text in enumerate(content.bullets[:9]):
                col = i % 3
                row = i // 3
                left = self.SUMMARY_COL_POSITIONS[col]
                top = self.SUMMARY_ROW_START + (row * self.SUMMARY_ROW_SPACING)
                text_box = slide.shapes.add_textbox(
                    left, top, self.SUMMARY_BOX_WIDTH, self.SUMMARY_BOX_HEIGHT
                )
                tf = text_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = bullet_text
                p.font.size = Pt(11)

    def _build_mailer_summary_slide(self, slide, content: SlideContent) -> None:
        """Build composite mailer summary slide -- 3 equal columns."""
        for ph in slide.placeholders:
            if ph.placeholder_format.idx != 0:
                try:
                    ph.element.getparent().remove(ph.element)
                except Exception:
                    pass

        COL1_L = Inches(0.2)
        COL2_L = Inches(4.4)
        COL3_L = Inches(8.8)
        COL_W = Inches(4.1)
        HEADER_SIZE = Pt(16)

        ROW1_TOP = Inches(1.6)
        KPI_VAL_TOP = Inches(2.1)
        KPI_LBL_TOP = Inches(2.55)
        SECT_TOP = Inches(3.2)
        CHART_TOP = Inches(3.5)

        # Title
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.38), Inches(9.0), Inches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content.title
        p.font.size = Pt(24)
        p.font.bold = False
        p.font.color.rgb = RGBColor(255, 255, 255)

        # Parse bullets
        insight_text = ""
        inside_numbers: list[str] = []
        if content.bullets:
            insight_text = content.bullets[0] if content.bullets[0] else ""
            inside_numbers = content.bullets[1:]

        # Insight text (upper-left)
        if insight_text:
            tb = slide.shapes.add_textbox(Inches(0.5), ROW1_TOP, Inches(5.0), Inches(1.0))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = insight_text
            p.font.size = Pt(15)
            p.font.color.rgb = RGBColor(0, 0, 0)

        # Mailer KPIs (upper-right)
        if content.kpis:
            kpi_left = Inches(7.8)
            kpi_block_w = Inches(5.2)
            tb = slide.shapes.add_textbox(kpi_left, ROW1_TOP, kpi_block_w, Inches(0.4))
            tf = tb.text_frame
            p = tf.paragraphs[0]
            p.text = "Mailer KPIs"
            p.font.size = HEADER_SIZE
            p.font.bold = True
            p.font.color.rgb = RGBColor(30, 61, 89)
            p.alignment = PP_ALIGN.CENTER

            kpi_items = list(content.kpis.items())
            kpi_each_w = 5.2 / max(len(kpi_items), 1)

            for i, (label_text, value) in enumerate(kpi_items):
                x = 7.8 + i * kpi_each_w

                tb = slide.shapes.add_textbox(
                    Inches(x), KPI_VAL_TOP, Inches(kpi_each_w), Inches(0.5)
                )
                tf = tb.text_frame
                p = tf.paragraphs[0]
                p.text = str(value)
                p.font.size = Pt(24)
                p.font.bold = True
                p.font.color.rgb = RGBColor(30, 61, 89)
                p.alignment = PP_ALIGN.CENTER

                tb = slide.shapes.add_textbox(
                    Inches(x), KPI_LBL_TOP, Inches(kpi_each_w), Inches(0.3)
                )
                tf = tb.text_frame
                p = tf.paragraphs[0]
                p.text = label_text
                p.font.size = Pt(13)
                p.font.color.rgb = RGBColor(0, 0, 0)
                p.alignment = PP_ALIGN.CENTER

        # Section headers (3 across)
        for col_left, header_text in [
            (COL1_L, "Response Share"),
            (COL2_L, "Response Rate"),
            (COL3_L, "Inside the Numbers"),
        ]:
            tb = slide.shapes.add_textbox(col_left, SECT_TOP, COL_W, Inches(0.3))
            tf = tb.text_frame
            p = tf.paragraphs[0]
            p.text = header_text
            p.font.size = HEADER_SIZE
            p.font.bold = True
            p.font.color.rgb = RGBColor(30, 61, 89)
            p.alignment = PP_ALIGN.CENTER

        # Donut chart (column 1)
        if content.images and len(content.images) > 0 and Path(content.images[0]).exists():
            slide.shapes.add_picture(content.images[0], COL1_L, CHART_TOP, width=COL_W)

        # Horizontal bar chart (column 2)
        if content.images and len(content.images) > 1 and Path(content.images[1]).exists():
            slide.shapes.add_picture(content.images[1], COL2_L, CHART_TOP, width=COL_W)

        # Inside the Numbers (column 3)
        if inside_numbers:
            # Dynamic spacing: compress to fit all items (up to 6)
            n_items = len(inside_numbers)
            row_h = min(1.2, 3.3 / max(n_items, 1))
            pct_size = Pt(22) if n_items > 4 else Pt(26)
            desc_size = Pt(11) if n_items > 4 else Pt(13)
            row_box_h = Inches(row_h * 0.85)

            for i, item in enumerate(inside_numbers):
                if "|" in item:
                    pct, desc = item.split("|", 1)
                else:
                    pct, desc = item, ""

                y_pos = 3.9 + i * row_h

                tb = slide.shapes.add_textbox(COL3_L, Inches(y_pos), Inches(1.4), row_box_h)
                tf = tb.text_frame
                tf.auto_size = None
                tf.word_wrap = False
                try:
                    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                    bodyPr = tf._txBody.find(qn("a:bodyPr"))
                    bodyPr.set("anchor", "ctr")
                except Exception:
                    pass
                p = tf.paragraphs[0]
                p.text = pct.strip()
                p.font.size = pct_size
                p.font.bold = True
                p.font.color.rgb = RGBColor(0, 80, 114)

                if desc:
                    tb = slide.shapes.add_textbox(
                        Inches(10.2), Inches(y_pos), Inches(2.6), row_box_h
                    )
                    tf = tb.text_frame
                    tf.auto_size = None
                    tf.word_wrap = True
                    try:
                        bodyPr = tf._txBody.find(qn("a:bodyPr"))
                        bodyPr.set("anchor", "ctr")
                    except Exception:
                        pass
                    p = tf.paragraphs[0]
                    p.text = desc.strip()
                    p.font.size = desc_size
                    p.font.color.rgb = RGBColor(0, 0, 0)


# =============================================================================
# SLIDE LAYOUT MAP -- slide_id -> (layout_index, slide_type)
# =============================================================================

SLIDE_LAYOUT_MAP: dict[str, tuple[int, str]] = {
    # DCTR -- penetration (data slides)
    "DCTR-1": (9, "screenshot"),
    "DCTR-2": (9, "screenshot"),
    "DCTR-3": (9, "screenshot"),
    "DCTR-4": (9, "screenshot"),
    "DCTR-5": (9, "screenshot"),
    "DCTR-6": (9, "screenshot"),
    "DCTR-7": (9, "screenshot"),
    "DCTR-8": (9, "screenshot"),
    "DCTR-9": (9, "screenshot"),
    "DCTR-10": (9, "screenshot"),
    "DCTR-11": (9, "screenshot"),
    "DCTR-12": (9, "screenshot"),
    "DCTR-13": (9, "screenshot"),
    "DCTR-14": (9, "screenshot"),
    "DCTR-15": (9, "screenshot"),
    "DCTR-16": (9, "screenshot"),
    # DCTR -- A7.x analysis charts
    "A7.4": (9, "screenshot"),
    "A7.5": (9, "screenshot"),
    "A7.6a": (9, "screenshot"),
    "A7.6b": (9, "screenshot"),
    "A7.7": (9, "screenshot"),
    "A7.8": (9, "screenshot"),
    "A7.9": (9, "screenshot"),
    "A7.10a": (13, "screenshot"),
    "A7.10b": (9, "screenshot"),
    "A7.10c": (4, "screenshot_kpi"),
    "A7.11": (9, "screenshot"),
    "A7.12": (9, "screenshot"),
    "A7.13": (9, "screenshot"),
    "A7.14": (9, "screenshot"),
    "A7.15": (9, "screenshot"),
    # Attrition
    "A9.1": (5, "screenshot_kpi"),
    "A9.2": (4, "screenshot"),
    "A9.3": (4, "screenshot"),
    "A9.4": (4, "screenshot"),
    "A9.5": (4, "screenshot"),
    "A9.6": (4, "screenshot"),
    "A9.7": (4, "screenshot"),
    "A9.8": (4, "screenshot"),
    "A9.9": (5, "screenshot_kpi"),
    "A9.10": (5, "screenshot_kpi"),
    "A9.11": (5, "screenshot_kpi"),
    "A9.12": (5, "screenshot_kpi"),
    "A9.13": (4, "screenshot"),
    # Reg E
    "A8.1": (9, "screenshot"),
    "A8.2": (9, "screenshot"),
    "A8.3": (9, "screenshot"),
    "A8.4a": (13, "screenshot"),
    "A8.4b": (9, "screenshot"),
    "A8.4c": (9, "screenshot"),
    "A8.5": (9, "screenshot"),
    "A8.6": (9, "screenshot"),
    "A8.7": (9, "screenshot"),
    "A8.10": (9, "screenshot"),
    "A8.11": (9, "screenshot"),
    "A8.12": (9, "screenshot"),
    "A8.13": (9, "screenshot"),
    # Value
    "A11.1": (13, "screenshot"),
    "A11.2": (13, "screenshot"),
    # Mailer
    "A13.5": (13, "screenshot"),
    "A13.6": (9, "screenshot"),
    "A14.2": (9, "screenshot"),
    "A15.1": (13, "screenshot"),
    "A15.2": (13, "screenshot"),
    "A15.3": (13, "screenshot"),
    "A15.4": (13, "screenshot"),
    # Mailer -- cohort trajectories
    "A16.1": (9, "screenshot"),
    "A16.2": (9, "screenshot"),
    "A16.3": (9, "screenshot"),
    "A16.4": (9, "screenshot"),
    "A16.5": (9, "screenshot"),
    "A16.6": (9, "screenshot"),
    # Mailer -- cumulative reach
    "A17.1": (9, "screenshot"),
    "A17.2": (9, "screenshot"),
    "A17.3": (9, "screenshot"),
    # Overview
    "A1": (9, "screenshot"),
    "A1b": (9, "screenshot"),
    "A3": (9, "screenshot"),
    # Insights
    "S1": (9, "screenshot"),
    "S2": (9, "screenshot"),
    "S3": (9, "screenshot"),
    "S4": (9, "screenshot"),
    "S5": (9, "screenshot"),
    "S6": (9, "screenshot"),
    "S7": (9, "screenshot"),
    "S8": (9, "screenshot"),
    # Insights -- effectiveness proof
    "A18.1": (9, "screenshot"),
    "A18.2": (9, "screenshot"),
    "A18.3": (9, "screenshot"),
    # Insights -- branch scorecard
    "A19.1": (13, "screenshot"),
    "A19.2": (9, "screenshot"),
    # Insights -- dormant opportunity
    "A20.1": (9, "screenshot"),
    "A20.2": (9, "screenshot"),
    "A20.3": (9, "screenshot"),
}


def _match_prefix(slide_id: str) -> tuple[int, str]:
    """Match slide_id by prefix for dynamic entries (e.g. A12.Nov25.Swipes)."""
    sid = slide_id.lower()
    if sid.startswith("a12"):
        return (13, "screenshot")
    if sid.startswith("a13") and sid not in ("a13.5", "a13.6"):
        return (13, "mailer_summary")
    if sid.startswith("a16"):
        return (9, "screenshot")
    if sid.startswith("a17"):
        return (9, "screenshot")
    if sid.startswith("a18"):
        return (9, "screenshot")
    if sid.startswith("a19"):
        return (9, "screenshot")
    if sid.startswith("a20"):
        return (9, "screenshot")
    return (9, "screenshot")


# =============================================================================
# CONSOLIDATION -- merge paired slides, separate appendix
# =============================================================================

DCTR_MERGES = [
    ("A7.6a", "A7.4", "DCTR Trajectory: Recent Trend & Segments"),
    ("A7.7", "A7.8", "DCTR Funnel: Historical vs TTM"),
    ("A7.11", "A7.12", "DCTR Opportunity: Age Analysis"),
]

DCTR_APPENDIX_IDS = {
    "A7.5",
    "A7.6b",
    "A7.13",
    "A7.14",
    "A7.15",
    "A7.9",
    "A7.10b",
    "A7.10c",
}

REGE_MERGES = [
    ("A8.10", "A8.11", "Reg E Funnel: All-Time vs TTM"),
    ("A8.5", "A8.6", "Reg E Opportunity: Age Analysis"),
]

REGE_APPENDIX_IDS = {
    "A8.7",
    "A8.4c",
    "A8.2",
    "A8.1",
    "A8.12",
    "A8.4b",
}

ATTRITION_MERGES = [
    ("A9.3", "A9.6", "Attrition Profile: Open vs Closed & Personal vs Business"),
]

ATTRITION_APPENDIX_IDS = {
    "A9.2",
    "A9.4",
    "A9.5",
    "A9.7",
    "A9.8",
    "A9.13",
}


def _consolidate(slides, merges, appendix_ids):
    """Merge paired slides and separate appendix slides.

    Returns (main_slides, appendix_slides).
    """
    merge_at = {}
    skip_ids = set()
    by_id = {r.slide_id: r for r in slides}

    for left_id, right_id, title in merges:
        left = by_id.get(left_id)
        right = by_id.get(right_id)
        if left and right:
            images = []
            if left.chart_path and left.chart_path.exists():
                images.append(str(left.chart_path))
            if right.chart_path and right.chart_path.exists():
                images.append(str(right.chart_path))

            merged_sc = SlideContent(
                slide_type="multi_screenshot",
                title=title,
                images=images,
                layout_index=6,
            )
            merge_at[left_id] = merged_sc
            skip_ids.add(left_id)
            skip_ids.add(right_id)

    result = []
    appendix_out = []
    for r in slides:
        sid = r.slide_id if hasattr(r, "slide_id") else getattr(r, "title", "")
        if sid in merge_at:
            result.append(merge_at[sid])
        elif sid in skip_ids:
            continue
        elif sid in appendix_ids:
            appendix_out.append(r)
        else:
            result.append(r)

    return result, appendix_out


# =============================================================================
# SECTION GROUPING
# =============================================================================

_SECTION_MAP = {
    "a1": "overview",
    "a1b": "overview",
    "a3": "overview",
    "a7": "dctr",
    "dctr": "dctr",
    "a8": "rege",
    "rege": "rege",
    "a9": "attrition",
    "att": "attrition",
    "a10": "value",
    "a11": "value",
    "val": "value",
    "a12": "mailer",
    "a13": "mailer",
    "a14": "mailer",
    "a15": "mailer",
    "a16": "mailer",
    "a17": "mailer",
    "mail": "mailer",
    "ics": "ics",
    "txn": "transaction",
    "m1": "transaction",
    "m2": "transaction",
    "m3": "transaction",
    "m4": "transaction",
    "m5": "transaction",
    "m6": "transaction",
    "m7": "transaction",
    "s1": "insights",
    "s2": "insights",
    "s3": "insights",
    "s4": "insights",
    "s5": "insights",
    "s6": "insights",
    "s7": "insights",
    "s8": "insights",
    "a18": "insights",
    "a19": "insights",
    "a20": "insights",
}

_SECTION_LABELS = {
    "overview": "Overview",
    "dctr": "Debit Card Transaction Revenue",
    "rege": "Regulation E",
    "attrition": "Attrition Analysis",
    "value": "Value Analysis",
    "mailer": "Mailer Analysis",
    "transaction": "Transaction Intelligence",
    "ics": "ICS Account Analysis",
    "insights": "Key Insights",
}


def _get_section(slide_id: str) -> str:
    """Get section name from slide_id."""
    prefix = slide_id.split("-")[0].split(".")[0].lower()
    return _SECTION_MAP.get(prefix, "other")


def _group_by_section(results: list) -> dict[str, list]:
    """Group AnalysisResult objects by section."""
    sections: dict[str, list] = {}
    for r in results:
        sid = getattr(r, "slide_id", "")
        section = _get_section(sid)
        sections.setdefault(section, []).append(r)
    return sections


# =============================================================================
# PREAMBLE BUILDER
# =============================================================================


def _build_preamble_slides(client_name: str, month: str) -> list[SlideContent]:
    """Build the 13 preamble slides that precede analysis content.

    These are title, agenda, section dividers, and blank placeholders
    for manual content (financial performance, revenue, lift matrix, etc.).
    """
    try:
        month_num = int(month.split(".")[1]) if "." in month else 1
        year = month.split(".")[0] if "." in month else "2026"
        month_name = calendar.month_name[month_num]
    except (ValueError, IndexError):
        month_name = ""
        year = ""

    title_date = f"{month_name} {year}" if month_name else month

    return [
        # P01: Intro
        SlideContent(
            slide_type="title",
            title=f"{client_name}\nAccount Revenue Solution | {title_date}",
            layout_index=1,
        ),
        # P02: Agenda (layout 13 = dark header)
        SlideContent(slide_type="blank", title="Agenda", layout_index=13),
        # P03: Program Performance divider
        SlideContent(
            slide_type="title",
            title=f"{client_name}\nProgram Performance | {title_date}",
            layout_index=1,
        ),
        # P04: Financial Performance -- blank for manual table
        SlideContent(
            slide_type="blank",
            title="Financial Performance",
            layout_index=0,
        ),
        # P05: Monthly Revenue -- blank
        SlideContent(
            slide_type="blank",
            title="Monthly Revenue \u2013 Last 12 Months",
            layout_index=12,
        ),
        # P06: ARS Lift Matrix -- blank
        SlideContent(
            slide_type="blank",
            title="ARS Lift Matrix",
            layout_index=8,
        ),
        # P07: ARS Mailer Revisit divider
        SlideContent(
            slide_type="title",
            title=f"{client_name}\nARS Mailer Revisit | {title_date}",
            layout_index=1,
        ),
        # P08: Swipes placeholder (will be wired to most recent A12 Swipes)
        SlideContent(
            slide_type="blank",
            title="ARS Mailer Revisit \u2013 Swipes",
            layout_index=13,
        ),
        # P09: Spend placeholder (will be wired to most recent A12 Spend)
        SlideContent(
            slide_type="blank",
            title="ARS Mailer Revisit \u2013 Spend",
            layout_index=13,
        ),
        # P10: DCO -- blank
        SlideContent(
            slide_type="blank",
            title="Data Check Overview\nOur goal is turning non-users and light-users into heavy users",
            layout_index=8,
        ),
        # P11: Mailer Summaries divider
        SlideContent(
            slide_type="title",
            title=f"Mailer Summaries\n{client_name} | {title_date}",
            layout_index=2,
        ),
        # P12: All Program Results -- blank
        SlideContent(
            slide_type="blank",
            title="All Program Results",
            layout_index=2,
        ),
        # P13: Program Responses to Date (will be wired to A13.5)
        SlideContent(
            slide_type="blank",
            title="Program Responses to Date",
            layout_index=13,
        ),
    ]


# =============================================================================
# RESULT -> SLIDE CONTENT CONVERSION
# =============================================================================


def _result_to_slide(result) -> SlideContent | None:
    """Convert an AnalysisResult to a SlideContent for the deck builder."""
    if not getattr(result, "success", True):
        return None

    chart_path = getattr(result, "chart_path", None)
    if not chart_path or not Path(chart_path).exists():
        return None

    slide_id = getattr(result, "slide_id", "")
    title = getattr(result, "title", "")
    kpis = getattr(result, "kpis", None)
    layout_idx = getattr(result, "layout_index", 5)
    slide_type = getattr(result, "slide_type", "screenshot")

    # Fall back to SLIDE_LAYOUT_MAP if module used defaults
    if layout_idx == 5 and slide_type == "screenshot":
        mapped = SLIDE_LAYOUT_MAP.get(slide_id)
        if mapped:
            layout_idx, slide_type = mapped
        else:
            layout_idx, slide_type = _match_prefix(slide_id)

    # Build image list (primary + extras)
    images = [str(chart_path)]
    extra = getattr(result, "extra_charts", None)
    if extra:
        images.extend(str(p) for p in extra if p and Path(p).exists())

    bullets = getattr(result, "bullets", None)

    return SlideContent(
        slide_type=slide_type,
        title=title,
        images=images,
        bullets=bullets,
        kpis=kpis,
        layout_index=layout_idx,
    )


# =============================================================================
# MAILER SLIDE ORDERING
# =============================================================================

_MAILER_NON_MONTH = {"A13.Agg", "A13.5", "A13.6", "A13", "A12"}
_MONTH_ABBRS = {
    "Jan": 1,
    "Feb": 2,
    "Mar": 3,
    "Apr": 4,
    "May": 5,
    "Jun": 6,
    "Jul": 7,
    "Aug": 8,
    "Sep": 9,
    "Oct": 10,
    "Nov": 11,
    "Dec": 12,
}


def _parse_mailer_month(slide_id: str) -> tuple[int, int] | None:
    """Extract (year, month_num) from slide IDs like A13.Jan26, A12.Feb26.Swipes."""
    import re

    m = re.search(r"\.([A-Z][a-z]{2})(\d{2})", slide_id)
    if m:
        abbr, yr = m.group(1), int(m.group(2))
        if abbr in _MONTH_ABBRS:
            return (2000 + yr, _MONTH_ABBRS[abbr])
    return None


def _consolidate_mailer(results: list) -> tuple[list, list]:
    """Split mailer results into main deck + appendix.

    Per-month groups: summary (A13.{month}) + swipes (A12.{month}.Swipes) + spend (A12.{month}.Spend).
    Most recent 2 months -> main. Older months -> appendix.
    A14.2 (mailer revisit) goes with most recent month.
    Aggregate and impact slides stay in main.
    """
    # Bucket slides
    month_slides: dict[tuple[int, int], list] = {}  # (year, month) -> [results]
    aggregate = []  # A13.Agg, A13.5
    revisit = []  # A14.x
    impact = []  # A15.x
    mailer_app = []  # A13.6 (rate trend) -> appendix
    other = []  # A12 or A13 without month suffix, etc.

    for r in results:
        sid = getattr(r, "slide_id", "")
        ym = _parse_mailer_month(sid)
        if ym:
            month_slides.setdefault(ym, []).append(r)
        elif sid.startswith("A13.Agg") or sid == "A13.5":
            aggregate.append(r)
        elif sid == "A13.6":
            mailer_app.append(r)
        elif sid.startswith("A14"):
            revisit.append(r)
        elif sid.startswith("A15"):
            impact.append(r)
        elif sid.startswith("A16"):
            impact.append(r)  # Cohort trajectories grouped with impact
        elif sid.startswith("A17"):
            impact.append(r)  # Cumulative reach grouped with impact
        else:
            other.append(r)

    # Sort months chronologically (most recent first)
    sorted_months = sorted(month_slides.keys(), reverse=True)

    # Within each month group, order: summary (A13) -> swipes (A12.Swipes) -> spend (A12.Spend) -> rest
    def _intra_month_key(r) -> int:
        sid = getattr(r, "slide_id", "")
        if sid.startswith("A13."):
            return 0  # summary first
        if "Swipes" in sid:
            return 1
        if "Spend" in sid:
            return 2
        return 3

    main_slides: list = []
    appendix_slides: list = []

    for i, ym in enumerate(sorted_months):
        group = sorted(month_slides[ym], key=_intra_month_key)
        if i < 2:
            # Most recent 2 months -> main
            main_slides.extend(group)
            if i == 0:
                # Mailer revisit goes after the most recent month
                main_slides.extend(revisit)
        else:
            appendix_slides.extend(group)

    # Aggregate summaries after monthly groups
    main_slides.extend(aggregate)
    # Impact slides last in main
    main_slides.extend(impact)
    # Other (catch-all)
    main_slides.extend(other)
    # Rate trend etc. to appendix
    appendix_slides.extend(mailer_app)

    return main_slides, appendix_slides


# =============================================================================
# BUILD DECK -- main entry point
# =============================================================================


def build_deck(ctx: PipelineContext) -> Path | None:
    """Build a PowerPoint deck from analysis results.

    Restores the full original ARS presentation flow:
    1. Preamble slides (13 intro/section/placeholder slides)
    2. Analysis slides grouped by section with proper layouts
    3. Consolidation (merge paired slides, separate appendix)
    4. Section ordering matching the reference deck
    """
    if not ctx.all_slides:
        logger.warning("No slides to build deck from")
        return None

    # Resolve template
    template = _FALLBACK_TEMPLATE
    if ctx.settings and hasattr(ctx.settings, "paths"):
        cfg_template = getattr(ctx.settings.paths, "template_path", None)
        if cfg_template and Path(cfg_template).exists():
            template = Path(cfg_template)

    if not template.exists():
        logger.warning("Template not found: {name}", name=template.name)
        return None

    _notify = ctx.progress_callback

    # Group results by section
    sections = _group_by_section(ctx.all_slides)

    dctr_results = sections.get("dctr", [])
    rege_results = sections.get("rege", [])
    attrition_results = sections.get("attrition", [])
    value_results = sections.get("value", [])
    mailer_results = sections.get("mailer", [])
    overview_results = sections.get("overview", [])
    insights_results = sections.get("insights", [])
    other_results = sections.get("other", [])

    if _notify:
        _notify("Building deck: consolidating slides...")

    # Consolidate: merge paired slides, separate appendix
    dctr_main, dctr_appendix = _consolidate(dctr_results, DCTR_MERGES, DCTR_APPENDIX_IDS)
    rege_main, rege_appendix = _consolidate(rege_results, REGE_MERGES, REGE_APPENDIX_IDS)
    attrition_main, attrition_appendix = _consolidate(
        attrition_results, ATTRITION_MERGES, ATTRITION_APPENDIX_IDS
    )

    # Separate value slides for DCTR and Reg E sections
    value_dctr = [r for r in value_results if getattr(r, "slide_id", "") == "A11.1"]
    value_rege = [r for r in value_results if getattr(r, "slide_id", "") == "A11.2"]

    # Build section subtitle
    client_name = ctx.client.client_name
    month = ctx.client.month
    try:
        month_num = int(month.split(".")[1]) if "." in month else 1
        year = month.split(".")[0] if "." in month else ""
        month_name = calendar.month_name[month_num]
        section_subtitle = f"{client_name} | {month_name} {year}"
    except (ValueError, IndexError):
        section_subtitle = client_name

    # Convert AnalysisResult -> SlideContent
    def _convert_list(items):
        converted = []
        for item in items:
            if isinstance(item, SlideContent):
                converted.append(item)
            else:
                sc = _result_to_slide(item)
                if sc:
                    converted.append(sc)
        return converted

    def _section_divider(title, subtitle=None, layout_index=1, slide_type="title"):
        full_title = f"{title}\n{subtitle}" if subtitle else title
        return SlideContent(slide_type=slide_type, title=full_title, layout_index=layout_index)

    # Build ordered analysis slides
    analysis_slides: list[SlideContent] = []

    # Mailer slides: per-month groups, most recent 2 in main, rest in appendix
    mailer_main, mailer_appendix = _consolidate_mailer(mailer_results)
    mailer_slides = _convert_list(mailer_main)
    mailer_app_slides = _convert_list(mailer_appendix)
    if mailer_slides:
        analysis_slides.extend(mailer_slides)

    # DCTR + Value of Debit Card
    dctr_slides = _convert_list(dctr_main)
    value_dctr_slides = _convert_list(value_dctr)
    if dctr_slides or value_dctr_slides:
        analysis_slides.append(_section_divider("Debit Card Take Rate", subtitle=section_subtitle))
        analysis_slides.extend(dctr_slides)
        analysis_slides.extend(value_dctr_slides)

    # Reg E + Value of Reg E
    rege_slides = _convert_list(rege_main)
    value_rege_slides = _convert_list(value_rege)
    if rege_slides or value_rege_slides:
        analysis_slides.append(_section_divider("Reg E Analysis", subtitle=section_subtitle))
        analysis_slides.extend(rege_slides)
        analysis_slides.extend(value_rege_slides)

    # Attrition
    attrition_slides = _convert_list(attrition_main)
    if attrition_slides:
        analysis_slides.append(_section_divider("Account Attrition", subtitle=section_subtitle))
        analysis_slides.extend(attrition_slides)

    # Summary placeholder
    analysis_slides.append(
        _section_divider("Summary & Key Takeaways", layout_index=12, slide_type="blank")
    )

    # Appendix
    dctr_app = _convert_list(dctr_appendix)
    rege_app = _convert_list(rege_appendix)
    attrition_app = _convert_list(attrition_appendix)
    overview_slides = _convert_list(overview_results)
    insights_slides = _convert_list(insights_results)
    other_slides = _convert_list(other_results)

    has_appendix = dctr_app or rege_app or attrition_app or overview_slides or mailer_app_slides
    if has_appendix:
        analysis_slides.append(_section_divider("Appendix", subtitle=section_subtitle))
        analysis_slides.extend(overview_slides)
        analysis_slides.extend(mailer_app_slides)
        analysis_slides.extend(dctr_app)
        analysis_slides.extend(rege_app)
        analysis_slides.extend(attrition_app)

    analysis_slides.extend(insights_slides)
    analysis_slides.extend(other_slides)

    # Build preamble
    preamble = _build_preamble_slides(client_name, month)

    # Wire preamble placeholders to actual results:
    # P08 (index 7) -> most recent A12.*.Swipes
    # P09 (index 8) -> most recent A12.*.Spend
    # P13 (index 12) -> A13.5 (count trend)
    _mailer_by_id = {getattr(r, "slide_id", ""): r for r in mailer_results}

    # Find most recent Swipes and Spend from A12 results
    _swipes = next(
        (
            _mailer_by_id[k]
            for k in sorted(_mailer_by_id, reverse=True)
            if k.startswith("A12.") and "swipe" in k.lower()
        ),
        None,
    )
    _spend = next(
        (
            _mailer_by_id[k]
            for k in sorted(_mailer_by_id, reverse=True)
            if k.startswith("A12.") and "spend" in k.lower()
        ),
        None,
    )
    _count_trend = _mailer_by_id.get("A13.5")

    for idx, result in [(7, _swipes), (8, _spend), (12, _count_trend)]:
        if result and idx < len(preamble):
            sc = _result_to_slide(result)
            if sc:
                preamble[idx] = sc

    # Combine
    final_slides = preamble + analysis_slides

    if not final_slides:
        logger.warning("No slides to build")
        return None

    if _notify:
        _notify(f"Building deck: {len(final_slides)} slides...")

    # Build the PPTX
    ctx.paths.pptx_dir.mkdir(parents=True, exist_ok=True)
    output_path = ctx.paths.pptx_dir / f"{ctx.client.client_id}_{ctx.client.month}_deck.pptx"

    try:
        builder = DeckBuilder(str(template))
        builder.build(final_slides, str(output_path))
        ctx.export_log.append(str(output_path))
        logger.info(
            "Deck built: {path} ({n} slides)",
            path=output_path.name,
            n=len(final_slides),
        )
        if _notify:
            _notify(f"Deck saved: {output_path.name} ({len(final_slides)} slides)")
        return output_path
    except Exception as exc:
        logger.error("Deck build failed: {err}", err=exc)
        if _notify:
            _notify(f"Deck build failed: {exc}")
        return None
