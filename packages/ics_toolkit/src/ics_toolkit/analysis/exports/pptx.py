"""PPTX report generator with data tables and chart images.

Produces two decks:
- **Primary** (~27 slides): Storyline-driven client-facing deck with merged
  charts, KPI panels, and ARS-style dark section dividers.
- **Secondary** (~70 slides): Full detail reference deck with table + chart
  for every analysis (the original 159-slide layout).
"""

import logging
from datetime import datetime
from io import BytesIO
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from ics_toolkit.analysis.analyses.base import AnalysisResult
from ics_toolkit.settings import AnalysisSettings as Settings

logger = logging.getLogger(__name__)

# -- Colors ----------------------------------------------------------------
NAVY = RGBColor(0x1B, 0x36, 0x5D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
ZEBRA_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
TOTAL_BG = RGBColor(0xE0, 0xE0, 0xE0)
LIGHT_GRAY = RGBColor(0x99, 0x99, 0x99)

# -- Named layout constants -- 2025-CSI-PPT-Template.pptx (20 layouts) -----
LAYOUT_TITLE_DARK = 0       # Title Slide (dark bg)
LAYOUT_TITLE = 1             # Title Slide_Reverse (light bg)
LAYOUT_SECTION = 4           # 2_Section Header
LAYOUT_SECTION_ALT = 5       # 5_Section Header (alt)
LAYOUT_CUSTOM = 8            # Custom Layout (wide title + open canvas)
LAYOUT_TWO_CONTENT = 9       # Two Content (side-by-side)
LAYOUT_BLANK = 11             # Blank (no placeholders)
LAYOUT_TITLE_RPE = 17         # 1_Title Slide_RPE
LAYOUT_TITLE_ICS = 19         # 5_Title Slide_ICS

# -- Table constraints -----------------------------------------------------
MAX_TABLE_ROWS = 12
MAX_TABLE_COLS = 10

# -- Slide positioning (widescreen 13.33" x 7.5") -------------------------
SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)

TITLE_LEFT = Inches(0.5)
TITLE_TOP = Inches(0.25)
TITLE_WIDTH = Inches(12.3)
TITLE_HEIGHT = Inches(0.55)

TABLE_LEFT = Inches(0.5)
TABLE_TOP = Inches(1.1)
TABLE_WIDTH = Inches(12.3)

CHART_LEFT = Inches(0.86)
CHART_TOP = Inches(1.5)
CHART_WIDTH = Inches(11.6)
CHART_MAX_HEIGHT = Inches(5.5)

# -- Merged slide positioning (two charts side-by-side) --------------------
MERGE_LEFT_X = Inches(0.86)
MERGE_RIGHT_X = Inches(6.81)
MERGE_IMG_Y = Inches(1.5)
MERGE_IMG_W = Inches(5.67)
MERGE_IMG_H = Inches(5.2)

# -- KPI panel positioning -------------------------------------------------
KPI_GRID_LEFT = Inches(0.5)
KPI_GRID_TOP = Inches(1.2)
KPI_BOX_W = Inches(3.8)
KPI_BOX_H = Inches(1.6)
KPI_COL_SPACING = Inches(4.1)
KPI_ROW_SPACING = Inches(1.8)

# Maps analysis name -> section grouping for slide organization.
SECTION_MAP = {
    "Executive Summary": [
        "Executive Summary",
    ],
    "Summary": [
        "Total ICS Accounts",
        "Open ICS Accounts",
        "ICS by Stat Code",
        "Product Code Distribution",
        "Debit Distribution",
        "Debit x Prod Code",
        "Debit x Branch",
        "ICS Penetration by Branch",
    ],
    "Portfolio Health": [
        "Net Portfolio Growth",
        "Engagement Decay",
        "Spend Concentration",
        "Closure by Source",
        "Closure by Branch",
        "Closure by Account Age",
        "Net Growth by Source",
        "Closure Rate Trend",
    ],
    "Source Analysis": [
        "Source Distribution",
        "Source x Stat Code",
        "Source x Prod Code",
        "Source x Branch",
        "Account Type",
        "Source by Year",
        "Source Acquisition Mix",
    ],
    "DM Source Deep-Dive": [
        "DM Overview",
        "DM by Branch",
        "DM by Debit Status",
        "DM by Product",
        "DM by Year Opened",
        "DM Activity Summary",
        "DM Activity by Branch",
        "DM Monthly Trends",
    ],
    "REF Source Deep-Dive": [
        "REF Overview",
        "REF by Branch",
        "REF by Debit Status",
        "REF by Product",
        "REF by Year Opened",
        "REF Activity Summary",
        "REF Activity by Branch",
        "REF Monthly Trends",
    ],
    "Demographics": [
        "Age Comparison",
        "Closures",
        "Open vs Close",
        "Balance Tiers",
        "Stat Open Close",
        "Age vs Balance",
        "Balance Tier Detail",
        "Age Distribution",
        "Balance Trajectory",
    ],
    "Activity Analysis": [
        "Activity Summary",
        "Activity by Debit+Source",
        "Activity by Balance",
        "Activity by Branch",
        "Monthly Trends",
        "Activity by Source Comparison",
        "Monthly Interchange Trend",
        "Business vs Personal",
    ],
    "Cohort Analysis": [
        "Cohort Activation",
        "Cohort Heatmap",
        "Cohort Milestones",
        "Activation Summary",
        "Growth Patterns",
        "Activation Personas",
        "Branch Activation",
    ],
    "Persona Deep-Dive": [
        "Persona Overview",
        "Persona Swipe Contribution",
        "Persona by Branch",
        "Persona by Source",
        "Persona Revenue Impact",
        "Persona by Balance Tier",
        "Persona Velocity",
        "Persona Cohort Trend",
    ],
    "Performance": [
        "Days to First Use",
        "Branch Performance Index",
        "Product Code Performance",
    ],
    "Strategic Insights": [
        "Activation Funnel",
        "Revenue Impact",
        "Revenue by Branch",
        "Revenue by Source",
        "Dormant High-Balance",
    ],
}

# All analysis names referenced in SECTION_MAP (for detecting un-mapped analyses).
_MAPPED_NAMES: set[str] = set()
for _names in SECTION_MAP.values():
    _MAPPED_NAMES.update(_names)

# Strings that mark a row as a "total" row.
_TOTAL_MARKERS = {"total", "grand total"}

# Pairs of analyses to merge side-by-side (left, right, merged title).
_SECTION_MERGE_PAIRS: list[tuple[str, str, str]] = [
    ("Closure by Branch", "Closure by Account Age", "Closures: Branch vs Account Age"),
]

# =========================================================================
# PRIMARY DECK: Storyline-driven layout
# =========================================================================

# Each section answers a business question.  Entries are either:
#   str             -> single chart-only slide
#   tuple[str, str] -> merged (side-by-side) slide from two analyses
# Analyses listed here appear in the Primary deck; everything else goes
# to the Secondary (detail/appendix) deck.

PRIMARY_STORYLINE: dict[str, list[str | tuple[str, str]]] = {
    "How Big Is This Program?": [
        "Executive Summary",
        "ICS Penetration by Branch",
        ("ICS by Stat Code", "Product Code Distribution"),
        "Activation Funnel",
    ],
    "Where Do ICS Accounts Come From?": [
        ("Source Distribution", "Source Acquisition Mix"),
        ("DM Overview", "REF Overview"),
        ("DM Activity Summary", "REF Activity Summary"),
    ],
    "How Engaged Are ICS Members?": [
        "Activity Summary",
        ("Monthly Trends", "Monthly Interchange Trend"),
        "Business vs Personal",
        "Spend Concentration",
    ],
    "How Fast Do Accounts Activate?": [
        "Cohort Heatmap",
        "Days to First Use",
        ("Persona Overview", "Persona Revenue Impact"),
        "Persona Swipe Contribution",
    ],
    "Which Branches Lead?": [
        "Branch Performance Index",
        ("Revenue by Branch", "Revenue by Source"),
    ],
    "What Are the Opportunities?": [
        "Dormant High-Balance",
        ("Closure by Source", "Closure by Branch"),
        "Closure Rate Trend",
        ("Age Comparison", "Balance Tiers"),
    ],
}

# Analyses whose hero_kpis metadata should render as a visual KPI panel
# instead of a data table in the primary deck.
KPI_PANEL_ANALYSES: set[str] = {
    "Executive Summary",
    "Activity Summary",
    "Dormant High-Balance",
    "DM Overview",
    "REF Overview",
    "DM Activity Summary",
    "REF Activity Summary",
}

# Flat set of all analysis names that appear in the primary storyline.
_PRIMARY_NAMES: set[str] = set()
for _entries in PRIMARY_STORYLINE.values():
    for _entry in _entries:
        if isinstance(_entry, tuple):
            _PRIMARY_NAMES.update(_entry)
        else:
            _PRIMARY_NAMES.add(_entry)


def _build_analysis_lookup(analyses: list[AnalysisResult]) -> dict[str, AnalysisResult]:
    """Build a name -> result lookup for successful analyses."""
    return {a.name: a for a in analyses if a.error is None}


# =========================================================================
# Public API
# =========================================================================


def write_chart_catalog(
    settings: Settings,
    analyses: list[AnalysisResult],
    chart_pngs: dict[str, bytes],
    output_path: Path | None = None,
) -> Path:
    """Build a chart catalog PPTX: one slide per chart with metadata labels.

    Each slide shows the chart index, name, rendered image, section,
    chart builder function name, and source .py file path.
    """
    import inspect

    from ics_toolkit.analysis.charts import CHART_REGISTRY

    if output_path is None:
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_path = settings.output_dir / f"Chart_Catalog_{timestamp}.pptx"

    output_path.parent.mkdir(parents=True, exist_ok=True)

    # Build chart name -> metadata mapping via introspection
    pkg_root = Path(__file__).resolve().parent.parent.parent
    chart_meta: dict[str, dict[str, str]] = {}
    for name, func in CHART_REGISTRY.items():
        try:
            src = Path(inspect.getfile(func)).resolve()
            rel = src.relative_to(pkg_root)
            source_str = f"ics_toolkit/{rel}"
        except (ValueError, TypeError):
            source_str = "unknown"
        chart_meta[name] = {"function": func.__name__, "source": source_str}

    # Build name -> section lookup
    section_lookup: dict[str, str] = {}
    for section_name, names in SECTION_MAP.items():
        for n in names:
            section_lookup[n] = section_name

    # Blank widescreen presentation (no template needed for catalog)
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, "ICS Chart Catalog")
    subtitle = (
        f"{settings.client_name or 'Client ' + (settings.client_id or 'unknown')}"
        f"  --  {len(chart_pngs)} charts"
    )
    txbox = slide.shapes.add_textbox(TITLE_LEFT, Inches(1.2), TITLE_WIDTH, Inches(1.0))
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_TEXT

    # One slide per chart
    for idx, (name, png_bytes) in enumerate(chart_pngs.items(), start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Title: "#idx  name"
        _add_slide_title(slide, f"#{idx}  {name}")

        # Chart image (reduced height to leave room for metadata)
        slide.shapes.add_picture(
            BytesIO(png_bytes),
            CHART_LEFT,
            CHART_TOP,
            width=CHART_WIDTH,
            height=Inches(4.5),
        )

        # Metadata text box
        meta = chart_meta.get(name, {})
        section = section_lookup.get(name, "Unmapped")
        func_name = meta.get("function", "N/A")
        source = meta.get("source", "N/A")

        meta_lines = [
            f"Section: {section}",
            f"Function: {func_name}()",
            f"Source: {source}",
        ]

        txbox = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(5.8),
            Inches(12.3),
            Inches(1.2),
        )
        tf = txbox.text_frame
        tf.word_wrap = True
        for i, line in enumerate(meta_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    prs.save(str(output_path))
    logger.info("Chart catalog saved: %s (%d slides)", output_path, len(prs.slides))
    return output_path


def write_ics_reports(
    settings: Settings,
    analyses: list[AnalysisResult],
    chart_pngs: dict[str, bytes] | None = None,
    output_dir: Path | None = None,
    per_section: bool = False,
) -> tuple[Path, Path] | tuple[Path, Path, dict[str, Path]]:
    """Generate Primary (storyline) and Secondary (detail) PPTX decks.

    Args:
        per_section: If True, also generate per-section module decks.

    Returns:
        (primary_path, secondary_path) or
        (primary_path, secondary_path, section_paths) when per_section=True.
    """
    out = output_dir or settings.output_dir
    out.mkdir(parents=True, exist_ok=True)

    date_str = datetime.now().strftime("%Y%m%d")
    client_id = settings.client_id or "unknown"

    primary_path = out / f"{client_id}_ICS_Primary_{date_str}.pptx"
    secondary_path = out / f"{client_id}_ICS_Secondary_{date_str}.pptx"

    primary = write_pptx_primary(
        settings,
        analyses,
        output_path=primary_path,
        chart_pngs=chart_pngs,
    )
    secondary = write_pptx_secondary(
        settings,
        analyses,
        output_path=secondary_path,
        chart_pngs=chart_pngs,
    )

    if per_section:
        section_dir = out / "sections"
        section_paths = write_pptx_per_section(
            settings,
            analyses,
            output_dir=section_dir,
            chart_pngs=chart_pngs,
        )
        return primary, secondary, section_paths

    return primary, secondary


def write_pptx_primary(
    settings: Settings,
    analyses: list[AnalysisResult],
    output_path: Path | None = None,
    chart_pngs: dict[str, bytes] | None = None,
) -> Path:
    """Build storyline-driven primary deck (~27 slides).

    - Charts only (no data tables -- those live in the Secondary deck).
    - Merge pairs render two charts side-by-side.
    - KPI-panel analyses render as visual KPI grids.
    - ARS-style dark section dividers.
    """
    if output_path is None:
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_path = settings.output_dir / f"ICS_Primary_{timestamp}.pptx"

    output_path.parent.mkdir(parents=True, exist_ok=True)

    lookup = _build_analysis_lookup(analyses)
    pngs = chart_pngs or {}

    prs = _create_presentation(settings)

    # Title slide
    _add_title_slide(prs, settings, title_text="ICS Analysis")

    for section_title, entries in PRIMARY_STORYLINE.items():
        # Check if any entry has data before adding divider
        has_content = False
        for entry in entries:
            if isinstance(entry, tuple):
                if entry[0] in lookup or entry[1] in lookup:
                    has_content = True
                    break
            elif entry in lookup:
                has_content = True
                break

        if not has_content:
            continue

        _add_styled_section_divider(prs, section_title)

        for entry in entries:
            if isinstance(entry, tuple):
                left_name, right_name = entry
                left_png = pngs.get(left_name)
                right_png = pngs.get(right_name)
                if left_png and right_png:
                    title = f"{left_name}  |  {right_name}"
                    _add_merged_slide(prs, title, left_png, right_png)
                elif left_png:
                    _add_chart_slide(prs, left_name, left_png)
                elif right_png:
                    _add_chart_slide(prs, right_name, right_png)
                else:
                    # Fall back to KPI panel or table for each
                    for name in (left_name, right_name):
                        if name in lookup:
                            _add_primary_analysis_slide(prs, name, lookup[name], pngs)
            else:
                if entry in lookup:
                    _add_primary_analysis_slide(prs, entry, lookup[entry], pngs)

    prs.save(str(output_path))
    logger.info(
        "Primary PPTX saved: %s (%d slides)",
        output_path,
        len(prs.slides),
    )
    return output_path


def write_pptx_secondary(
    settings: Settings,
    analyses: list[AnalysisResult],
    output_path: Path | None = None,
    chart_pngs: dict[str, bytes] | None = None,
) -> Path:
    """Build full detail/appendix deck (table + chart for every analysis).

    This is the original write_pptx_report behaviour with all 80 analyses.
    """
    if output_path is None:
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_path = settings.output_dir / f"ICS_Secondary_{timestamp}.pptx"

    output_path.parent.mkdir(parents=True, exist_ok=True)

    lookup = _build_analysis_lookup(analyses)
    pngs = chart_pngs or {}

    prs = _create_presentation(settings)

    # Title slide
    _add_title_slide(prs, settings, title_text="ICS Analysis -- Detail Appendix")

    # Walk SECTION_MAP in order
    for section_name, analysis_names in SECTION_MAP.items():
        section_analyses = [(name, lookup[name]) for name in analysis_names if name in lookup]
        if not section_analyses:
            continue

        _add_styled_section_divider(prs, section_name)

        # Pairs to merge side-by-side instead of separate slides
        _merged: set[str] = set()
        for left_name, right_name, _merge_title in _SECTION_MERGE_PAIRS:
            if left_name in pngs and right_name in pngs:
                _merged.update({left_name, right_name})

        for name, analysis in section_analyses:
            _add_table_slide(prs, analysis)

            # Emit merged chart slide for paired analyses
            if name in _merged:
                for left_name, right_name, merge_title in _SECTION_MERGE_PAIRS:
                    if name == left_name and left_name in pngs and right_name in pngs:
                        _add_merged_slide(prs, merge_title, pngs[left_name], pngs[right_name])
                continue

            if name in pngs:
                _add_chart_slide(prs, analysis.title, pngs[name])

    # Catch unmapped analyses
    unmapped = [(a.name, a) for a in analyses if a.error is None and a.name not in _MAPPED_NAMES]
    if unmapped:
        _add_section_divider(prs, "Additional Analyses")
        for name, analysis in unmapped:
            _add_table_slide(prs, analysis)
            if name in pngs:
                _add_chart_slide(prs, analysis.title, pngs[name])

    prs.save(str(output_path))
    logger.info(
        "Secondary PPTX saved: %s (%d slides)",
        output_path,
        len(prs.slides),
    )
    return output_path


def write_pptx_per_section(
    settings: Settings,
    analyses: list[AnalysisResult],
    output_dir: Path | None = None,
    chart_pngs: dict[str, bytes] | None = None,
    sections: list[str] | None = None,
) -> dict[str, Path]:
    """Generate one PPTX deck per SECTION_MAP section.

    Each deck contains a title slide, styled section divider, and the
    table + chart slides for that section's analyses only.

    Args:
        settings: Analysis settings.
        analyses: List of analysis results.
        output_dir: Directory for output files.
        chart_pngs: Optional dict of chart PNG bytes keyed by analysis name.
        sections: Optional list of section names to generate.
            If None, generates all sections that have data.

    Returns:
        Dict mapping section name to output file path.
    """
    out = output_dir or settings.output_dir
    out.mkdir(parents=True, exist_ok=True)

    date_str = datetime.now().strftime("%Y%m%d")
    client_id = settings.client_id or "unknown"

    lookup = _build_analysis_lookup(analyses)
    pngs = chart_pngs or {}

    target_sections = sections or list(SECTION_MAP.keys())
    results: dict[str, Path] = {}

    for section_name in target_sections:
        if section_name not in SECTION_MAP:
            logger.warning("Unknown section: %s (skipping)", section_name)
            continue

        analysis_names = SECTION_MAP[section_name]
        section_analyses = [(name, lookup[name]) for name in analysis_names if name in lookup]
        if not section_analyses:
            continue

        # Sanitize section name for filename
        safe_name = section_name.replace(" ", "_").replace("/", "-")
        path = out / f"{client_id}_ICS_{safe_name}_{date_str}.pptx"

        prs = _create_presentation(settings)
        _add_title_slide(prs, settings, title_text=f"ICS Analysis -- {section_name}")
        _add_styled_section_divider(prs, section_name)

        # Merge pairs (same logic as secondary deck)
        _merged: set[str] = set()
        for left_name, right_name, _merge_title in _SECTION_MERGE_PAIRS:
            if left_name in pngs and right_name in pngs:
                _merged.update({left_name, right_name})

        for name, analysis in section_analyses:
            _add_table_slide(prs, analysis)

            if name in _merged:
                for left_name, right_name, merge_title in _SECTION_MERGE_PAIRS:
                    if name == left_name and left_name in pngs and right_name in pngs:
                        _add_merged_slide(prs, merge_title, pngs[left_name], pngs[right_name])
                continue

            if name in pngs:
                _add_chart_slide(prs, analysis.title, pngs[name])

        prs.save(str(path))
        results[section_name] = path
        logger.info(
            "Section deck saved: %s (%d slides)",
            path.name,
            len(prs.slides),
        )

    logger.info("Generated %d section decks", len(results))
    return results


def write_pptx_report(
    settings: Settings,
    analyses: list[AnalysisResult],
    output_path: Path | None = None,
    chart_pngs: dict[str, bytes] | None = None,
) -> Path:
    """Build and save a PPTX presentation with tables and charts.

    Every successful analysis gets a data-table slide.
    Analyses with chart PNGs get an additional chart slide.

    .. deprecated::
        Use :func:`write_ics_reports` for Primary + Secondary decks,
        or :func:`write_pptx_secondary` for the full detail deck.
    """
    return write_pptx_secondary(
        settings,
        analyses,
        output_path=output_path,
        chart_pngs=chart_pngs,
    )


# =========================================================================
# Presentation setup
# =========================================================================


def _create_presentation(settings: Settings) -> Presentation:
    """Load 2025-CSI-PPT-Template or create blank widescreen presentation."""
    from pptx.oxml.ns import qn

    template = settings.pptx_template
    if template and Path(template).exists():
        prs = Presentation(str(template))
        # Remove sample slides shipped with 2025 template
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].get(qn("r:id"))
            prs.part.drop_rel(rId)
            prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])
        logger.info("Loaded template: %s", template)
    else:
        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
        if template:
            logger.warning("Template not found: %s (using blank)", template)
    return prs


def _get_layout(prs: Presentation, preferred: int, fallback: int = LAYOUT_CUSTOM):
    """Get slide layout by index with fallback."""
    for idx in [preferred, fallback, LAYOUT_CUSTOM, LAYOUT_BLANK, 0]:
        try:
            return prs.slide_layouts[idx]
        except IndexError:
            continue
    return prs.slide_layouts[0]


# =========================================================================
# Slide builders
# =========================================================================


def _add_title_slide(
    prs: Presentation,
    settings: Settings,
    title_text: str = "ICS Accounts Analysis",
) -> None:
    """Add branded title slide."""
    layout = _get_layout(prs, preferred=LAYOUT_TITLE_ICS, fallback=LAYOUT_TITLE)
    slide = prs.slides.add_slide(layout)

    if slide.shapes.title:
        slide.shapes.title.text = title_text

    subtitle = settings.client_name or f"Client {settings.client_id}"
    for ph_idx in [1, 13, 14]:
        try:
            slide.placeholders[ph_idx].text = subtitle
            break
        except (KeyError, IndexError):
            continue


def _add_section_divider(prs: Presentation, title: str) -> None:
    """Add section divider slide (template-based, simple)."""
    layout = _get_layout(prs, preferred=LAYOUT_SECTION, fallback=LAYOUT_SECTION_ALT)
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = title


def _add_styled_section_divider(prs: Presentation, title: str) -> None:
    """Add ARS-style dark section divider with navy background and white text."""
    layout = _get_layout(prs, preferred=LAYOUT_SECTION, fallback=LAYOUT_CUSTOM)
    slide = prs.slides.add_slide(layout)

    # Set the slide background to navy
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = NAVY

    # Clear any existing placeholder text
    if slide.shapes.title:
        slide.shapes.title.text = ""

    # Add large centered title
    txbox = slide.shapes.add_textbox(
        Inches(1.5),
        Inches(2.5),
        Inches(10.33),
        Inches(2.0),
    )
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Decorative line below title
    p2 = tf.add_paragraph()
    p2.text = "\u2500" * 30
    p2.font.size = Pt(14)
    p2.font.color.rgb = LIGHT_GRAY
    p2.alignment = PP_ALIGN.CENTER


def _add_merged_slide(
    prs: Presentation,
    title: str,
    left_png: bytes,
    right_png: bytes,
) -> None:
    """Add slide with two chart images side by side."""
    layout = _get_layout(prs, preferred=LAYOUT_TWO_CONTENT, fallback=LAYOUT_CUSTOM)
    slide = prs.slides.add_slide(layout)

    _add_slide_title(slide, title)

    slide.shapes.add_picture(
        BytesIO(left_png),
        MERGE_LEFT_X,
        MERGE_IMG_Y,
        width=MERGE_IMG_W,
        height=MERGE_IMG_H,
    )
    slide.shapes.add_picture(
        BytesIO(right_png),
        MERGE_RIGHT_X,
        MERGE_IMG_Y,
        width=MERGE_IMG_W,
        height=MERGE_IMG_H,
    )


def _add_kpi_slide(prs: Presentation, analysis: AnalysisResult) -> None:
    """Add a KPI panel slide from hero_kpis metadata (2x3 grid)."""
    hero_kpis = analysis.metadata.get("hero_kpis", {}) if analysis.metadata else {}
    if not hero_kpis:
        return

    layout = _get_layout(prs, preferred=LAYOUT_CUSTOM, fallback=LAYOUT_BLANK)
    slide = prs.slides.add_slide(layout)
    _add_slide_title(slide, analysis.title)

    traffic_lights = analysis.metadata.get("traffic_lights", {}) if analysis.metadata else {}

    items = list(hero_kpis.items())
    for idx, (key, value) in enumerate(items[:6]):
        col = idx % 3
        row = idx // 3

        left = KPI_GRID_LEFT + Inches(col * 4.1)
        top = KPI_GRID_TOP + Inches(row * 2.5)

        # KPI value box
        txbox = slide.shapes.add_textbox(left, top, KPI_BOX_W, Inches(0.6))
        tf = txbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = _format_kpi_display(key, value)
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p.alignment = PP_ALIGN.CENTER

        # KPI label
        label_box = slide.shapes.add_textbox(left, top + Inches(0.7), KPI_BOX_W, Inches(0.4))
        lf = label_box.text_frame
        lf.word_wrap = True
        lp = lf.paragraphs[0]
        lp.text = key
        lp.font.size = Pt(12)
        lp.font.color.rgb = DARK_TEXT
        lp.alignment = PP_ALIGN.CENTER

        # Traffic light indicator
        light = traffic_lights.get(key, "")
        if light and light != "gray":
            ind_box = slide.shapes.add_textbox(
                left,
                top + Inches(1.1),
                KPI_BOX_W,
                Inches(0.3),
            )
            ip = ind_box.text_frame.paragraphs[0]
            ip.text = f"[{light.upper()}]"
            ip.font.size = Pt(10)
            ip.font.color.rgb = _traffic_light_color(light)
            ip.alignment = PP_ALIGN.CENTER


def _add_primary_analysis_slide(
    prs: Presentation,
    name: str,
    analysis: AnalysisResult,
    pngs: dict[str, bytes],
) -> None:
    """Add the best available slide for a single analysis in the primary deck.

    Priority: chart PNG > KPI panel > skip (no tables in primary).
    """
    if name in pngs:
        _add_chart_slide(prs, analysis.title, pngs[name])
    elif name in KPI_PANEL_ANALYSES:
        _add_kpi_slide(prs, analysis)
    # No table fallback in primary deck -- tables go to secondary only


def _add_table_slide(prs: Presentation, analysis: AnalysisResult) -> None:
    """Add slide with analysis title and formatted data table."""
    df = analysis.df
    if df.empty:
        return

    # Show most-recent month first for temporal tables
    if "Month" in df.columns:
        total_mask = df["Month"].astype(str).str.strip().str.lower().isin(_TOTAL_MARKERS)
        total_rows = df[total_mask]
        data_rows = df[~total_mask].iloc[::-1]
        df = pd.concat([data_rows, total_rows], ignore_index=True)

    layout = _get_layout(prs, preferred=LAYOUT_CUSTOM, fallback=LAYOUT_BLANK)
    slide = prs.slides.add_slide(layout)

    # Title
    _add_slide_title(slide, analysis.title)

    # Prepare data for table (truncate if needed)
    show_df, truncated = _prepare_table_df(df)
    cols = list(show_df.columns[:MAX_TABLE_COLS])
    col_truncated = len(df.columns) > MAX_TABLE_COLS

    nrows = len(show_df) + 1  # +1 for header
    ncols = len(cols)

    # Calculate table height: ~0.45" per row (Pt 20 font)
    row_height = 0.45
    table_height = Inches(nrows * row_height)

    tbl_shape = slide.shapes.add_table(
        nrows,
        ncols,
        TABLE_LEFT,
        TABLE_TOP,
        TABLE_WIDTH,
        table_height,
    )
    table = tbl_shape.table

    # Distribute column widths
    col_width = int(TABLE_WIDTH / ncols)
    for j in range(ncols):
        table.columns[j].width = col_width

    # Header row
    for j, col_name in enumerate(cols):
        cell = table.cell(0, j)
        cell.text = str(col_name)
        _style_header_cell(cell)

    # Data rows
    for i, (_, row) in enumerate(show_df.iterrows()):
        is_total = _is_total_row(row)
        is_odd = i % 2 == 1

        for j, col_name in enumerate(cols):
            cell = table.cell(i + 1, j)
            cell.text = _format_cell_value(row[col_name], col_name)
            _style_data_cell(cell, is_total=is_total, is_odd=is_odd)

    # Truncation note
    if truncated or col_truncated:
        parts = []
        if truncated:
            parts.append(f"Showing {len(show_df)} of {len(df)} rows")
        if col_truncated:
            parts.append(f"{len(cols)} of {len(df.columns)} columns")
        note = " | ".join(parts) + " -- see Excel report for full data"
        note_top = TABLE_TOP + table_height + Inches(0.15)
        _add_footnote(slide, note, note_top)


def _add_chart_slide(prs: Presentation, title: str, png_bytes: bytes) -> None:
    """Add slide with chart image."""
    layout = _get_layout(prs, preferred=LAYOUT_CUSTOM, fallback=LAYOUT_BLANK)
    slide = prs.slides.add_slide(layout)

    _add_slide_title(slide, title)

    slide.shapes.add_picture(
        BytesIO(png_bytes),
        CHART_LEFT,
        CHART_TOP,
        width=CHART_WIDTH,
        height=CHART_MAX_HEIGHT,
    )


# =========================================================================
# Shared slide helpers
# =========================================================================


def _add_slide_title(slide, text: str) -> None:
    """Add a title text box to the slide."""
    txbox = slide.shapes.add_textbox(
        TITLE_LEFT,
        TITLE_TOP,
        TITLE_WIDTH,
        TITLE_HEIGHT,
    )
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = NAVY


def _add_footnote(slide, text: str, top) -> None:
    """Add a small footnote text box."""
    txbox = slide.shapes.add_textbox(
        TABLE_LEFT,
        top,
        TABLE_WIDTH,
        Inches(0.3),
    )
    p = txbox.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(8)
    p.font.italic = True
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)


# =========================================================================
# KPI formatting helpers
# =========================================================================


def _format_kpi_display(key: str, value: object) -> str:
    """Format a KPI value for large-font display on a slide."""
    if isinstance(value, float) and ("Rate" in key or "%" in key):
        return f"{value:.1f}%"
    if isinstance(value, (int, float)) and any(
        kw in key for kw in ("Interchange", "Revenue", "Balance", "Spend")
    ):
        return f"${value:,.0f}"
    if isinstance(value, (int, float)):
        return f"{value:,}"
    return str(value)


def _traffic_light_color(light: str) -> RGBColor:
    """Map a traffic-light string to an RGB color."""
    _map = {
        "green": RGBColor(0x22, 0x8B, 0x22),
        "yellow": RGBColor(0xCC, 0xAA, 0x00),
        "red": RGBColor(0xCC, 0x33, 0x33),
    }
    return _map.get(light.lower(), DARK_TEXT)


# =========================================================================
# Table helpers
# =========================================================================


def _prepare_table_df(df: pd.DataFrame) -> tuple[pd.DataFrame, bool]:
    """Truncate DataFrame for slide display, preserving total rows.

    Returns (display_df, was_truncated).
    """
    if len(df) <= MAX_TABLE_ROWS:
        return df, False

    # Preserve total/grand total rows at bottom
    total_mask = df.iloc[:, 0].astype(str).str.strip().str.lower().isin(_TOTAL_MARKERS)
    total_rows = df[total_mask]

    if len(total_rows) > 0:
        n_total = len(total_rows)
        data_rows = df[~total_mask].head(MAX_TABLE_ROWS - n_total)
        return pd.concat([data_rows, total_rows]), True

    return df.head(MAX_TABLE_ROWS), True


def _is_total_row(row: pd.Series) -> bool:
    """Detect if a row is a Total/Grand Total summary row."""
    first = str(row.iloc[0]).strip().lower()
    return first in _TOTAL_MARKERS


def _format_cell_value(val, col_name: str) -> str:
    """Format a cell value for slide display."""
    if pd.isna(val):
        return ""

    col_lower = col_name.lower()

    # Percentage columns
    if "%" in col_name or "rate" in col_lower or "pct" in col_lower:
        if isinstance(val, (int, float)):
            return f"{val:.1f}%"

    # Currency columns
    if any(kw in col_lower for kw in ("balance", "spend", "revenue", "interchange")):
        if isinstance(val, (int, float)):
            if abs(val) >= 1000:
                return f"${val:,.0f}"
            return f"${val:,.2f}"

    # Integer counts
    if isinstance(val, (int,)):
        return f"{val:,}"

    # Floats
    if isinstance(val, float):
        if val == int(val) and abs(val) < 1e9:
            return f"{int(val):,}"
        return f"{val:,.1f}"

    return str(val)


def _style_header_cell(cell) -> None:
    """Style a table header cell: navy background, white bold text."""
    cell.fill.solid()
    cell.fill.fore_color.rgb = NAVY

    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.bold = True
        paragraph.font.color.rgb = WHITE
        paragraph.alignment = PP_ALIGN.CENTER


def _style_data_cell(
    cell,
    is_total: bool = False,
    is_odd: bool = False,
) -> None:
    """Style a data cell with alternating shading and total highlighting."""
    if is_total:
        cell.fill.solid()
        cell.fill.fore_color.rgb = TOTAL_BG
    elif is_odd:
        cell.fill.solid()
        cell.fill.fore_color.rgb = ZEBRA_GRAY
    else:
        cell.fill.background()

    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = DARK_TEXT
        paragraph.alignment = PP_ALIGN.CENTER
        if is_total:
            paragraph.font.bold = True
