"""Referral intelligence PPTX report -- self-contained slide builder."""

from __future__ import annotations

import logging
from io import BytesIO
from pathlib import Path

from ics_toolkit.analysis.analyses.base import AnalysisResult
from ics_toolkit.settings import ReferralSettings

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    Presentation = None  # type: ignore[assignment,misc]

logger = logging.getLogger(__name__)

SLIDE_WIDTH = Inches(13.33) if Presentation else None
SLIDE_HEIGHT = Inches(7.5) if Presentation else None

REFERRAL_SECTION_MAP = {
    "Referral Intelligence Overview": ["Overview KPIs"],
    "Referrer Influence": [
        "Top Referrers",
        "Emerging Referrers",
        "Dormant High-Value Referrers",
        "One-time vs Repeat Referrers",
    ],
    "Staff & Branch": [
        "Staff Multipliers",
        "Branch Influence Density",
    ],
    "Code Health": ["Code Health Report"],
}

_MAPPED_NAMES: set[str] = set()
for _names in REFERRAL_SECTION_MAP.values():
    _MAPPED_NAMES.update(_names)


def _build_analysis_lookup(analyses: list[AnalysisResult]) -> dict[str, AnalysisResult]:
    """Build a name -> result lookup for successful analyses."""
    return {a.name: a for a in analyses if a.error is None}


def _get_layout(prs, preferred: int = 1, fallback: int = 0):
    """Safely get a slide layout by index."""
    try:
        return prs.slide_layouts[preferred]
    except (IndexError, KeyError):
        return prs.slide_layouts[fallback]


def _add_section_divider(prs, section_name: str) -> None:
    """Add a section divider slide."""
    layout = _get_layout(prs, preferred=2, fallback=0)
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = section_name


def _add_chart_slide(prs, title: str, png_bytes: bytes) -> None:
    """Add a slide with a chart image."""
    blank = _get_layout(prs, preferred=6, fallback=0)
    slide = prs.slides.add_slide(blank)

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
    txBox.text_frame.paragraphs[0].text = title
    txBox.text_frame.paragraphs[0].font.size = Pt(18)
    txBox.text_frame.paragraphs[0].font.bold = True

    slide.shapes.add_picture(
        BytesIO(png_bytes),
        Inches(1.5),
        Inches(1.2),
        Inches(10),
        Inches(5.5),
    )


def _add_table_slide(prs, analysis: AnalysisResult) -> None:
    """Add a slide with a data table from an AnalysisResult."""
    df = analysis.df
    if df is None or df.empty:
        return

    blank = _get_layout(prs, preferred=6, fallback=0)
    slide = prs.slides.add_slide(blank)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
    txBox.text_frame.paragraphs[0].text = analysis.title
    txBox.text_frame.paragraphs[0].font.size = Pt(18)
    txBox.text_frame.paragraphs[0].font.bold = True

    # Table (limit rows for readability)
    max_rows = min(len(df), 20)
    rows = max_rows + 1  # +1 for header
    cols = len(df.columns)
    table_shape = slide.shapes.add_table(
        rows, cols, Inches(0.5), Inches(1.2), Inches(12), Inches(5.5)
    )
    table = table_shape.table

    # Header row
    for j, col_name in enumerate(df.columns):
        cell = table.cell(0, j)
        cell.text = str(col_name)

    # Data rows
    for i in range(max_rows):
        for j in range(cols):
            cell = table.cell(i + 1, j)
            val = df.iloc[i, j]
            cell.text = str(val) if val is not None else ""


def write_referral_pptx(
    settings: ReferralSettings,
    analyses: list[AnalysisResult],
    output_path: Path | None = None,
    chart_pngs: dict[str, bytes] | None = None,
) -> Path:
    """Build and save a PPTX presentation for referral intelligence."""
    if output_path is None:
        output_path = settings.output_dir / "Referral_Intelligence.pptx"

    output_path.parent.mkdir(parents=True, exist_ok=True)

    lookup = _build_analysis_lookup(analyses)
    pngs = chart_pngs or {}

    # Create presentation
    template = settings.pptx_template
    if template and Path(template).exists():
        prs = Presentation(str(template))
        logger.info("Loaded template: %s", template)
    else:
        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
        if template:
            logger.warning("Template not found: %s (using blank)", template)

    # Title slide
    layout = _get_layout(prs, preferred=1, fallback=0)
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = "Referral Intelligence Report"
    subtitle = settings.client_name or (f"Client {settings.client_id or 'unknown'}")
    for ph_idx in [1, 13, 14]:
        try:
            slide.placeholders[ph_idx].text = subtitle
            break
        except (KeyError, IndexError):
            continue

    # Walk section map
    for section_name, analysis_names in REFERRAL_SECTION_MAP.items():
        section_analyses = [(name, lookup[name]) for name in analysis_names if name in lookup]
        if not section_analyses:
            continue

        _add_section_divider(prs, section_name)

        for name, analysis in section_analyses:
            _add_table_slide(prs, analysis)
            if name in pngs:
                _add_chart_slide(prs, analysis.title, pngs[name])

    # Catch unmapped analyses
    unmapped = [(a.name, a) for a in analyses if a.error is None and a.name not in _MAPPED_NAMES]
    if unmapped:
        _add_section_divider(prs, "Additional Analyses")
        for name, analysis in unmapped:
            _add_table_slide(prs, analysis)
            if name in pngs:
                _add_chart_slide(prs, analysis.title, pngs[name])

    prs.save(str(output_path))
    logger.info(
        "Referral PPTX report saved: %s (%d slides)",
        output_path,
        len(prs.slides),
    )
    return output_path
