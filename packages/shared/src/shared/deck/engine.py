"""Shared PPTX deck engine -- DeckBuilder + SlideContent.

Extracted from the ICS and ARS deck builders to provide a single reusable
engine for assembling SlideContent objects into PowerPoint presentations.
"""

from __future__ import annotations

import logging
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

logger = logging.getLogger(__name__)

_FALLBACK_TEMPLATE = Path(__file__).parent / "template" / "2025-CSI-PPT-Template.pptx"

MAX_CHART_HEIGHT = Inches(5.0)

# Named layout constants -- 2025-CSI-PPT-Template.pptx (20 layouts)
LAYOUT_TITLE_DARK = 0
LAYOUT_TITLE = 1
LAYOUT_SECTION = 4
LAYOUT_SECTION_ALT = 5
LAYOUT_CUSTOM = 8
LAYOUT_TWO_CONTENT = 9
LAYOUT_BLANK = 11
LAYOUT_TITLE_RPE = 17
LAYOUT_TITLE_ARS = 18
LAYOUT_TITLE_ICS = 19


@dataclass
class SlideContent:
    """Container for all information needed to build a single slide.

    Attributes
    ----------
    slide_type:
        Determines which builder method creates this slide.
        Options: 'title', 'section', 'screenshot', 'screenshot_kpi',
                 'multi_screenshot', 'summary', 'blank'.
    title:
        Text displayed in the slide's title area.
    images:
        File paths to chart/screenshot PNG images.
    kpis:
        Key-value pairs for KPI callouts. For title slides use
        ``{"subtitle": "January 2025"}``.
    bullets:
        Text items for summary slides (up to 9, displayed in 3x3 grid).
    layout_index:
        Index of slide layout from the template to use.
    """

    slide_type: str
    title: str
    images: list[str] | None = None
    kpis: dict[str, str] | None = None
    bullets: list[str] | None = None
    layout_index: int = 8  # LAYOUT_CUSTOM (2025-CSI-PPT-Template)


class DeckBuilder:
    """Assembles SlideContent objects into a PowerPoint presentation.

    Parameters
    ----------
    template_path:
        Path to a .pptx template file.  If ``None``, a blank presentation
        is created instead.
    """

    # -- Single screenshot defaults (inches) -- 2025-CSI-PPT-Template ---------
    SINGLE_IMG_LEFT = Inches(0.86)
    SINGLE_IMG_TOP = Inches(1.6)
    SINGLE_IMG_WIDTH = Inches(11.6)

    SINGLE_IMG_RIGHT_LEFT = Inches(5.5)
    SINGLE_IMG_RIGHT_TOP = Inches(1.6)
    SINGLE_IMG_RIGHT_WIDTH = Inches(7.0)

    # -- Multi screenshot - standard layout -----------------------------------
    MULTI_IMG_STD_TOP = Inches(1.82)
    MULTI_IMG_STD_LEFT_POS = Inches(0.86)
    MULTI_IMG_STD_RIGHT_POS = Inches(6.81)
    MULTI_IMG_STD_WIDTH = Inches(5.67)

    # -- Screenshot with KPIs ------------------------------------------------
    KPI_IMG_LEFT = Inches(0.5)
    KPI_IMG_TOP = Inches(1.5)
    KPI_IMG_WIDTH = Inches(6.5)

    KPI_TEXT_LEFT = Inches(7.2)
    KPI_TEXT_TOP_START = Inches(1.8)
    KPI_TEXT_WIDTH = Inches(5.5)
    KPI_VALUE_HEIGHT = Inches(0.5)
    KPI_LABEL_HEIGHT = Inches(0.3)
    KPI_SPACING = Inches(1.0)

    # -- Summary slide (3x3 bullet grid) --------------------------------------
    SUMMARY_COL_POSITIONS = [Inches(0.86), Inches(4.86), Inches(8.86)]
    SUMMARY_ROW_START = Inches(1.8)
    SUMMARY_ROW_SPACING = Inches(1.2)
    SUMMARY_BOX_WIDTH = Inches(3.5)
    SUMMARY_BOX_HEIGHT = Inches(1.0)

    _BUILDER_REGISTRY: dict[str, str] = {
        "title": "_build_title_slide",
        "section": "_build_section_slide",
        "screenshot": "_build_screenshot_slide",
        "screenshot_kpi": "_build_screenshot_kpi_slide",
        "multi_screenshot": "_build_multi_screenshot_slide",
        "summary": "_build_summary_slide",
        "blank": "_build_blank_slide",
    }

    def __init__(self, template_path: Path | str | None = None) -> None:
        if template_path is not None:
            self.template_path: Path | None = Path(template_path)
        else:
            self.template_path = None
        self.prs: Presentation | None = None

    # -- positioning helpers --------------------------------------------------

    def _get_single_positioning(self, layout_index: int) -> tuple:
        """Return (left, top, width) in EMU for a single-chart layout."""
        # LAYOUT_CUSTOM (8) -- workhorse layout
        if layout_index == LAYOUT_CUSTOM:
            return (Inches(0.86), Inches(1.6), Inches(11.6))
        # LAYOUT_BLANK (11) -- full canvas
        if layout_index == LAYOUT_BLANK:
            return (Inches(0.86), Inches(1.6), Inches(11.6))
        # LAYOUT_TWO_CONTENT (9) -- single chart in left half
        if layout_index == LAYOUT_TWO_CONTENT:
            return (Inches(0.86), Inches(1.82), Inches(5.67))
        # Default
        return (Inches(0.86), Inches(1.6), Inches(11.6))

    def _get_multi_positioning(self, layout_index: int) -> tuple:
        """Return (top, left_pos, right_pos, width) in EMU for multi-chart layouts."""
        return (Inches(1.82), Inches(0.86), Inches(6.81), Inches(5.67))

    # -- build ----------------------------------------------------------------

    def build(self, slides: list[SlideContent], output_path: Path | str) -> Path:
        """Build complete PowerPoint deck from slide definitions."""
        from pptx.oxml.ns import qn

        output_path = Path(output_path)

        if self.template_path is not None:
            self.prs = Presentation(str(self.template_path))
            # Remove sample slides shipped with 2025 template
            while len(self.prs.slides) > 0:
                rId = self.prs.slides._sldIdLst[0].get(qn("r:id"))
                self.prs.part.drop_rel(rId)
                self.prs.slides._sldIdLst.remove(self.prs.slides._sldIdLst[0])
        else:
            self.prs = Presentation()

        n_layouts = len(self.prs.slide_layouts)

        for i, slide_content in enumerate(slides):
            if slide_content.layout_index >= n_layouts:
                logger.warning(
                    "Slide %d '%s' layout_index=%d but template has %d layouts, using 0",
                    i,
                    slide_content.title[:40],
                    slide_content.layout_index,
                    n_layouts,
                )
                slide_content.layout_index = 0
            try:
                self._add_slide(slide_content)
            except Exception as exc:
                logger.error(
                    "Slide %d '%s' (type=%s, layout=%d) failed: %s",
                    i,
                    slide_content.title[:40],
                    slide_content.slide_type,
                    slide_content.layout_index,
                    exc,
                )

        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(output_path))
        logger.info("Deck saved: %s (%d slides)", output_path, len(slides))
        return output_path

    def _add_slide(self, content: SlideContent) -> None:
        """Create a slide and dispatch to the appropriate builder method."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[content.layout_index])

        method_name = self._BUILDER_REGISTRY.get(content.slide_type)
        if method_name is None:
            logger.warning("Unknown slide_type: '%s'", content.slide_type)
            return

        builder = getattr(self, method_name)
        builder(slide, content)

    # -- image helper ---------------------------------------------------------

    @staticmethod
    def _add_fitted_picture(slide, img_path, left, top, max_width, max_height=None):
        """Add image scaled to fit within max_width and max_height."""
        effective_max_h = max_height or MAX_CHART_HEIGHT
        try:
            import warnings

            from PIL import Image

            with warnings.catch_warnings():
                warnings.filterwarnings("ignore", category=Image.DecompressionBombWarning)
                with Image.open(img_path) as img:
                    native_w, native_h = img.size
            aspect = native_h / native_w
            height_at_width = int(max_width * aspect)
            if height_at_width > effective_max_h:
                slide.shapes.add_picture(img_path, left, top, height=effective_max_h)
            else:
                slide.shapes.add_picture(img_path, left, top, width=max_width)
        except ImportError:
            slide.shapes.add_picture(img_path, left, top, width=max_width)

    # -- shared title helper --------------------------------------------------

    @staticmethod
    def _set_title_and_subtitle(slide, content: SlideContent) -> None:
        """Parse 'Title\\nSubtitle' and set placeholders."""
        title_text = content.title
        subtitle_text = None

        if "\n" in content.title:
            parts = content.title.split("\n", 1)
            title_text = parts[0]
            subtitle_text = parts[1] if len(parts) > 1 else None

        if slide.shapes.title:
            slide.shapes.title.text = title_text

        if subtitle_text:
            try:
                slide.placeholders[13].text = subtitle_text
            except (KeyError, IndexError):
                pass

    # -- title slide ----------------------------------------------------------

    def _build_title_slide(self, slide, content: SlideContent) -> None:
        """Build title slide with main title and optional subtitle."""
        if content.layout_index == 1:
            self._build_title_layout_1(slide, content)
            return

        title_lines = content.title.split("\n") if "\n" in content.title else [content.title]

        if slide.shapes.title:
            slide.shapes.title.text = title_lines[0]

        additional_text: list[str] = title_lines[1:]
        if content.kpis and "subtitle" in content.kpis:
            additional_text.append(content.kpis["subtitle"])

        text_placeholders = [26, 29, 30, 27, 28, 31]
        for i, text in enumerate(additional_text):
            if i >= len(text_placeholders):
                break
            try:
                slide.placeholders[text_placeholders[i]].text = text
            except (KeyError, IndexError):
                try:
                    slide.placeholders[1].text = text
                except (KeyError, IndexError):
                    pass

    def _build_title_layout_1(self, slide, content: SlideContent) -> None:
        """Build layout-1 title slide with a custom text box."""
        title_lines = content.title.split("\n") if "\n" in content.title else [content.title]
        subtitle = content.kpis.get("subtitle", "") if content.kpis else ""

        full_text = title_lines[0]
        if len(title_lines) > 1:
            full_text += f"\n{title_lines[1]}"
        if len(title_lines) > 2:
            full_text += f"\n{title_lines[2]}"
        elif subtitle:
            separator = "\u2500" * 20
            full_text += f"\n{separator}\n{subtitle}"

        text_box = slide.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(6.0), Inches(2.0))
        tf = text_box.text_frame
        tf.word_wrap = True

        lines = full_text.split("\n")
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.alignment = PP_ALIGN.LEFT
            if i == 0:
                p.font.size = Pt(32)
                p.font.bold = True
            else:
                p.font.size = Pt(18)

    # -- section slide --------------------------------------------------------

    def _build_section_slide(self, slide, content: SlideContent) -> None:
        """Build section divider slide with title and optional subtitle."""
        title_text = content.title
        subtitle_text = None

        if "\n" in content.title:
            parts = content.title.split("\n", 1)
            title_text = parts[0]
            subtitle_text = parts[1] if len(parts) > 1 else None

        if slide.shapes.title:
            slide.shapes.title.text = title_text

        if subtitle_text:
            for ph_idx in [44, 13, 1, 14]:
                try:
                    slide.placeholders[ph_idx].text = subtitle_text
                    break
                except (KeyError, IndexError):
                    continue

    # -- screenshot slide -----------------------------------------------------

    def _build_screenshot_slide(self, slide, content: SlideContent) -> None:
        """Build slide with title, optional subtitle, and a single image."""
        self._set_title_and_subtitle(slide, content)

        left, top, width = self._get_single_positioning(content.layout_index)

        if content.images and Path(content.images[0]).exists():
            self._add_fitted_picture(slide, content.images[0], left, top, width)
        elif content.images:
            logger.warning("Image not found, skipping: %s", content.images[0])

    # -- screenshot_kpi slide -------------------------------------------------

    def _build_screenshot_kpi_slide(self, slide, content: SlideContent) -> None:
        """Build slide with image on left and KPI callouts on right."""
        self._set_title_and_subtitle(slide, content)

        left, top, width = self._get_single_positioning(content.layout_index)

        if content.images and Path(content.images[0]).exists():
            self._add_fitted_picture(slide, content.images[0], left, top, width)
        elif content.images:
            logger.warning("Image not found, skipping: %s", content.images[0])

        if not content.kpis:
            return

        kpi_placeholder_pairs = [
            (26, 19),
            (27, 28),
        ]

        kpi_items = [(k, v) for k, v in content.kpis.items() if k != "subtitle"]

        for i, (label, value) in enumerate(kpi_items):
            if i >= len(kpi_placeholder_pairs):
                break

            label_idx, value_idx = kpi_placeholder_pairs[i]

            try:
                slide.placeholders[label_idx].text = label
            except (KeyError, IndexError):
                pass

            try:
                slide.placeholders[value_idx].text = str(value)
            except (KeyError, IndexError):
                pass

    # -- multi_screenshot slide -----------------------------------------------

    def _build_multi_screenshot_slide(self, slide, content: SlideContent) -> None:
        """Build slide with two images side by side."""
        self._set_title_and_subtitle(slide, content)

        top, left_pos, right_pos, width = self._get_multi_positioning(content.layout_index)

        positions = [
            (left_pos, top, width),
            (right_pos, top, width),
        ]

        if not content.images:
            return

        for i, img_path in enumerate(content.images[:2]):
            if Path(img_path).exists():
                img_left, img_top, img_width = positions[i]
                self._add_fitted_picture(slide, img_path, img_left, img_top, img_width)
            else:
                logger.warning("Image not found, skipping: %s", img_path)

    # -- summary slide --------------------------------------------------------

    def _build_summary_slide(self, slide, content: SlideContent) -> None:
        """Build summary slide with bullets in a 3x3 grid."""
        if slide.shapes.title:
            slide.shapes.title.text = content.title

        if not content.bullets:
            return

        max_bullets = 9
        for i, bullet in enumerate(content.bullets[:max_bullets]):
            col = i % 3
            row = i // 3

            left = self.SUMMARY_COL_POSITIONS[col]
            top = self.SUMMARY_ROW_START + (row * self.SUMMARY_ROW_SPACING)

            txbox = slide.shapes.add_textbox(
                left, top, self.SUMMARY_BOX_WIDTH, self.SUMMARY_BOX_HEIGHT
            )
            tf = txbox.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = bullet
            tf.paragraphs[0].font.size = Pt(14)

    # -- blank slide ----------------------------------------------------------

    def _build_blank_slide(self, slide, content: SlideContent) -> None:
        """Build blank placeholder slide."""
        if slide.shapes.title:
            slide.shapes.title.text = content.title
