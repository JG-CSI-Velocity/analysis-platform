"""PowerPoint report generator for TXN analysis (charts-only deck)."""

from __future__ import annotations

import logging
from datetime import datetime
from io import BytesIO
from pathlib import Path
from typing import TYPE_CHECKING

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

if TYPE_CHECKING:
    from txn_analysis.pipeline import PipelineResult

logger = logging.getLogger(__name__)

SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)

TITLE_LEFT = Inches(0.5)
TITLE_TOP = Inches(0.25)
TITLE_WIDTH = Inches(12.3)
TITLE_HEIGHT = Inches(0.6)

IMG_LEFT = Inches(0.6)
IMG_TOP = Inches(1.1)
IMG_WIDTH = Inches(12.1)
IMG_HEIGHT = Inches(5.8)

NAVY = RGBColor(0x1B, 0x36, 0x5D)


def _add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = title
    else:
        slide.shapes.add_textbox(TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT).text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle
    else:
        sub_box = slide.shapes.add_textbox(
            TITLE_LEFT, TITLE_TOP + Inches(0.7), TITLE_WIDTH, Inches(0.4)
        )
        sub_box.text_frame.text = subtitle


def _add_chart_slide(prs: Presentation, title: str, img_bytes: bytes | None) -> None:
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)

    title_box = slide.shapes.add_textbox(TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT)
    title_tf = title_box.text_frame
    title_tf.text = title
    title_tf.paragraphs[0].font.size = Pt(24)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = NAVY

    if img_bytes:
        slide.shapes.add_picture(
            BytesIO(img_bytes), IMG_LEFT, IMG_TOP, width=IMG_WIDTH, height=IMG_HEIGHT
        )


def write_pptx_report(result: PipelineResult, path: Path, chart_pngs: dict[str, bytes]) -> None:
    """Generate a simple charts-only PPTX deck."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    client = result.settings.client_name or result.settings.client_id or "Client"
    date_str = datetime.now().strftime("%B %d, %Y")
    _add_title_slide(prs, "Transaction Analysis", f"{client} • {date_str}")

    if not chart_pngs:
        _add_chart_slide(prs, "No charts available", None)
        prs.save(str(path))
        return

    charts_by_analysis: dict[str, list[tuple[str, bytes]]] = {}
    for key, png in chart_pngs.items():
        analysis_key = key.split(":")[0]
        charts_by_analysis.setdefault(analysis_key, []).append((key, png))

    analyses_by_name = {a.name: a for a in result.analyses}
    seen: set[str] = set()

    for analysis in result.analyses:
        charts = charts_by_analysis.get(analysis.name, [])
        if not charts:
            continue
        seen.add(analysis.name)
        for key, png in charts:
            suffix = ""
            if ":" in key:
                suffix = f" — {key.split(':', 1)[1].replace('_', ' ').title()}"
            title = (analysis.title or analysis.name.replace("_", " ").title()) + suffix
            _add_chart_slide(prs, title, png)

    # Append charts that didn't match an analysis name
    for analysis_key, charts in charts_by_analysis.items():
        if analysis_key in seen:
            continue
        for key, png in charts:
            title = analysis_key.replace("_", " ").title()
            if ":" in key:
                title += f" — {key.split(':', 1)[1].replace('_', ' ').title()}"
            _add_chart_slide(prs, title, png)

    prs.save(str(path))
    logger.info("PPTX report: %s", path)
