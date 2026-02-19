"""Tests for the PowerPoint deck builder."""

from pathlib import Path

import pandas as pd
from pptx import Presentation

from ars_analysis.analytics.base import AnalysisResult
from ars_analysis.output.deck_builder import _group_by_section, build_deck
from ars_analysis.pipeline.context import ClientInfo, OutputPaths, PipelineContext


def _make_ctx(tmp_path, with_chart=False):
    """Build a PipelineContext with test slides."""
    ctx = PipelineContext(
        client=ClientInfo(client_id="1234", client_name="Test CU", month="2026.01"),
        paths=OutputPaths.from_base(tmp_path, "1234", "2026.01"),
    )

    chart_path = None
    if with_chart:
        # Create a minimal PNG file
        ctx.paths.charts_dir.mkdir(parents=True, exist_ok=True)
        chart_path = ctx.paths.charts_dir / "test_chart.png"
        _write_minimal_png(chart_path)

    ctx.all_slides = [
        AnalysisResult(
            slide_id="DCTR-1",
            title="Overall DCTR",
            chart_path=chart_path,
            notes="Overall debit card penetration rate.",
        ),
        AnalysisResult(
            slide_id="A9.1",
            title="Attrition Rate",
            excel_data={"summary": pd.DataFrame({"Rate": [0.05]})},
            notes="Annual attrition rate.",
        ),
    ]
    return ctx


def _write_minimal_png(path: Path):
    """Write a minimal valid PNG file (1x1 red pixel)."""
    # Minimal PNG: 8-byte signature + IHDR + IDAT + IEND
    import struct
    import zlib

    signature = b"\x89PNG\r\n\x1a\n"

    # IHDR chunk: width=1, height=1, bit_depth=8, color_type=2 (RGB)
    ihdr_data = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
    ihdr_crc = zlib.crc32(b"IHDR" + ihdr_data) & 0xFFFFFFFF
    ihdr = struct.pack(">I", 13) + b"IHDR" + ihdr_data + struct.pack(">I", ihdr_crc)

    # IDAT chunk: filter byte (0) + RGB pixel (255, 0, 0)
    raw_data = zlib.compress(b"\x00\xff\x00\x00")
    idat_crc = zlib.crc32(b"IDAT" + raw_data) & 0xFFFFFFFF
    idat = struct.pack(">I", len(raw_data)) + b"IDAT" + raw_data + struct.pack(">I", idat_crc)

    # IEND chunk
    iend_crc = zlib.crc32(b"IEND") & 0xFFFFFFFF
    iend = struct.pack(">I", 0) + b"IEND" + struct.pack(">I", iend_crc)

    path.write_bytes(signature + ihdr + idat + iend)


class TestGroupBySection:
    """_group_by_section maps slide IDs to sections."""

    def test_dctr_prefix(self):
        slides = [AnalysisResult(slide_id="DCTR-1", title="Test")]
        groups = _group_by_section(slides)
        assert "dctr" in groups
        assert len(groups["dctr"]) == 1

    def test_a7_prefix(self):
        slides = [AnalysisResult(slide_id="A7.4", title="Test")]
        groups = _group_by_section(slides)
        assert "dctr" in groups

    def test_a9_prefix(self):
        slides = [AnalysisResult(slide_id="A9.1", title="Test")]
        groups = _group_by_section(slides)
        assert "attrition" in groups

    def test_ics_prefix(self):
        slides = [AnalysisResult(slide_id="ICS-1", title="Test")]
        groups = _group_by_section(slides)
        assert "ics" in groups

    def test_multiple_sections(self):
        slides = [
            AnalysisResult(slide_id="DCTR-1", title="DCTR Test"),
            AnalysisResult(slide_id="A9.1", title="Attrition Test"),
            AnalysisResult(slide_id="ICS-1", title="ICS Test"),
        ]
        groups = _group_by_section(slides)
        assert len(groups) == 3

    def test_unknown_prefix(self):
        slides = [AnalysisResult(slide_id="XYZ-1", title="Unknown")]
        groups = _group_by_section(slides)
        assert "xyz" in groups


class TestBuildDeck:
    """build_deck produces a valid PowerPoint file."""

    def test_no_slides_returns_none(self, tmp_path):
        ctx = PipelineContext(
            client=ClientInfo(client_id="1234", client_name="Test CU", month="2026.01"),
            paths=OutputPaths.from_base(tmp_path, "1234", "2026.01"),
        )
        result = build_deck(ctx)
        assert result is None

    def test_creates_pptx_with_chart(self, tmp_path):
        ctx = _make_ctx(tmp_path, with_chart=True)
        result = build_deck(ctx)
        assert result is not None
        assert result.exists()
        assert result.suffix == ".pptx"

    def test_pptx_has_slides(self, tmp_path):
        ctx = _make_ctx(tmp_path, with_chart=True)
        result = build_deck(ctx)
        prs = Presentation(str(result))
        # At least: 1 divider + 1 chart slide for dctr, 1 divider + 1 text slide for attrition
        assert len(prs.slides) >= 4

    def test_export_log_updated(self, tmp_path):
        ctx = _make_ctx(tmp_path, with_chart=True)
        build_deck(ctx)
        assert len(ctx.export_log) == 1
        assert "deck.pptx" in ctx.export_log[0]

    def test_text_only_slides(self, tmp_path):
        ctx = _make_ctx(tmp_path, with_chart=False)
        result = build_deck(ctx)
        assert result is not None
        prs = Presentation(str(result))
        # Divider slides + text slides
        assert len(prs.slides) >= 2

    def test_output_filename(self, tmp_path):
        ctx = _make_ctx(tmp_path, with_chart=True)
        result = build_deck(ctx)
        assert result.name == "1234_2026.01_deck.pptx"

    def test_failed_slides_skipped(self, tmp_path):
        ctx = _make_ctx(tmp_path, with_chart=True)
        ctx.all_slides.append(
            AnalysisResult(slide_id="DCTR-99", title="Failed", success=False, error="Test error")
        )
        result = build_deck(ctx)
        prs = Presentation(str(result))
        # Failed slide should not be added
        slide_titles = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    slide_titles.append(shape.text_frame.text)
        assert "Failed" not in slide_titles
