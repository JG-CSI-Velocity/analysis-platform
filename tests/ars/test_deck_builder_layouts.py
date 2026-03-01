"""Tests for deck builder layout mapping, consolidation, preamble, and slide conversion."""

from pathlib import Path

import pytest
from pptx import Presentation

from ars_analysis.analytics.base import AnalysisResult
from ars_analysis.output.deck_builder import (
    ATTRITION_APPENDIX_IDS,
    ATTRITION_MERGES,
    DCTR_APPENDIX_IDS,
    DCTR_MERGES,
    LAYOUT_CUSTOM,
    LAYOUT_SECTION,
    LAYOUT_TITLE_RPE,
    SLIDE_LAYOUT_MAP,
    DeckBuilder,
    SlideContent,
    _build_preamble_slides,
    _consolidate,
    _consolidate_mailer,
    _get_section,
    _group_by_section,
    _match_prefix,
    _result_to_slide,
    build_deck,
)
from ars_analysis.pipeline.context import ClientInfo, OutputPaths, PipelineContext


def _write_minimal_png(path: Path):
    """Write a minimal valid PNG file (1x1 red pixel)."""
    import struct
    import zlib

    signature = b"\x89PNG\r\n\x1a\n"
    ihdr_data = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
    ihdr_crc = zlib.crc32(b"IHDR" + ihdr_data) & 0xFFFFFFFF
    ihdr = struct.pack(">I", 13) + b"IHDR" + ihdr_data + struct.pack(">I", ihdr_crc)
    raw_data = zlib.compress(b"\x00\xff\x00\x00")
    idat_crc = zlib.crc32(b"IDAT" + raw_data) & 0xFFFFFFFF
    idat = struct.pack(">I", len(raw_data)) + b"IDAT" + raw_data + struct.pack(">I", idat_crc)
    iend_crc = zlib.crc32(b"IEND") & 0xFFFFFFFF
    iend = struct.pack(">I", 0) + b"IEND" + struct.pack(">I", iend_crc)
    path.write_bytes(signature + ihdr + idat + iend)


# =============================================================================
# SLIDE_LAYOUT_MAP
# =============================================================================


class TestSlideLayoutMap:
    """SLIDE_LAYOUT_MAP has correct structure and valid values."""

    def test_map_not_empty(self):
        assert len(SLIDE_LAYOUT_MAP) > 50

    def test_all_layout_indices_in_range(self):
        for slide_id, (layout_idx, _) in SLIDE_LAYOUT_MAP.items():
            assert 0 <= layout_idx <= 19, f"{slide_id} has out-of-range layout {layout_idx}"

    def test_all_slide_types_valid(self):
        valid = {"screenshot", "screenshot_kpi", "multi_screenshot", "mailer_summary", "blank"}
        for slide_id, (_, slide_type) in SLIDE_LAYOUT_MAP.items():
            assert slide_type in valid, f"{slide_id} has invalid type {slide_type}"

    def test_dctr_entries_present(self):
        dctr_ids = [f"DCTR-{i}" for i in range(1, 17)]
        for sid in dctr_ids:
            assert sid in SLIDE_LAYOUT_MAP, f"Missing {sid}"

    def test_a7_entries_present(self):
        expected = [
            "A7.4",
            "A7.5",
            "A7.6a",
            "A7.6b",
            "A7.7",
            "A7.8",
            "A7.9",
            "A7.10a",
            "A7.10b",
            "A7.10c",
            "A7.11",
            "A7.12",
            "A7.13",
            "A7.14",
            "A7.15",
        ]
        for sid in expected:
            assert sid in SLIDE_LAYOUT_MAP

    def test_attrition_entries_present(self):
        expected = [f"A9.{i}" for i in range(1, 14)]
        for sid in expected:
            assert sid in SLIDE_LAYOUT_MAP

    def test_rege_entries_present(self):
        expected = [
            "A8.1",
            "A8.2",
            "A8.3",
            "A8.4a",
            "A8.4b",
            "A8.4c",
            "A8.5",
            "A8.6",
            "A8.7",
            "A8.10",
            "A8.11",
            "A8.12",
            "A8.13",
        ]
        for sid in expected:
            assert sid in SLIDE_LAYOUT_MAP


# =============================================================================
# _match_prefix
# =============================================================================


class TestMatchPrefix:
    """Prefix matching for dynamic slide IDs."""

    def test_a12_prefix(self):
        assert _match_prefix("A12.Nov25.Swipes") == (LAYOUT_CUSTOM, "screenshot")

    def test_a13_monthly_mailer_summary(self):
        # A13.{month} uses mailer_summary (3-column: donut + hbar + inside numbers)
        assert _match_prefix("A13.Sep24") == (LAYOUT_CUSTOM, "mailer_summary")

    def test_a13_5_not_mailer(self):
        # A13.5 and A13.6 are excluded from the monthly prefix match
        layout, stype = _match_prefix("a13.5")
        assert stype == "screenshot"  # falls through to default

    def test_unknown_prefix_default(self):
        assert _match_prefix("ZZZ.1") == (LAYOUT_CUSTOM, "screenshot")


# =============================================================================
# _mailer_sort_key
# =============================================================================


class TestConsolidateMailer:
    """Mailer consolidation: per-month groups, most recent 2 in main, rest in appendix."""

    def test_monthly_groups_chronological(self):
        results = [
            AnalysisResult(slide_id="A13.Jan26", title="Jan26 Summary"),
            AnalysisResult(slide_id="A12.Jan26.Swipes", title="Jan26 Swipes"),
            AnalysisResult(slide_id="A12.Jan26.Spend", title="Jan26 Spend"),
            AnalysisResult(slide_id="A13.Feb26", title="Feb26 Summary"),
            AnalysisResult(slide_id="A12.Feb26.Swipes", title="Feb26 Swipes"),
            AnalysisResult(slide_id="A12.Feb26.Spend", title="Feb26 Spend"),
        ]
        main, appendix = _consolidate_mailer(results)
        ids = [r.slide_id for r in main]
        # Feb26 (most recent) first, then Jan26
        assert ids[0] == "A13.Feb26"
        assert ids[1] == "A12.Feb26.Swipes"
        assert ids[2] == "A12.Feb26.Spend"
        assert ids[3] == "A13.Jan26"

    def test_most_recent_2_in_main(self):
        results = [
            AnalysisResult(slide_id="A13.Nov25", title="Nov"),
            AnalysisResult(slide_id="A13.Dec25", title="Dec"),
            AnalysisResult(slide_id="A13.Jan26", title="Jan"),
        ]
        main, appendix = _consolidate_mailer(results)
        main_ids = {r.slide_id for r in main}
        app_ids = {r.slide_id for r in appendix}
        assert "A13.Jan26" in main_ids
        assert "A13.Dec25" in main_ids
        assert "A13.Nov25" in app_ids

    def test_revisit_after_most_recent(self):
        results = [
            AnalysisResult(slide_id="A14.2", title="Revisit"),
            AnalysisResult(slide_id="A13.Jan26", title="Jan"),
            AnalysisResult(slide_id="A13.Feb26", title="Feb"),
        ]
        main, appendix = _consolidate_mailer(results)
        ids = [r.slide_id for r in main]
        feb_idx = ids.index("A13.Feb26")
        rev_idx = ids.index("A14.2")
        jan_idx = ids.index("A13.Jan26")
        # Revisit should be after Feb26 group but before Jan26
        assert rev_idx > feb_idx
        assert rev_idx < jan_idx

    def test_aggregate_and_impact_in_main(self):
        results = [
            AnalysisResult(slide_id="A13.Agg", title="Aggregate"),
            AnalysisResult(slide_id="A13.5", title="Count Trend"),
            AnalysisResult(slide_id="A15.1", title="Impact"),
        ]
        main, appendix = _consolidate_mailer(results)
        main_ids = {r.slide_id for r in main}
        assert "A13.Agg" in main_ids
        assert "A13.5" in main_ids
        assert "A15.1" in main_ids
        assert len(appendix) == 0

    def test_intra_month_order(self):
        results = [
            AnalysisResult(slide_id="A12.Jan26.Spend", title="Spend"),
            AnalysisResult(slide_id="A13.Jan26", title="Summary"),
            AnalysisResult(slide_id="A12.Jan26.Swipes", title="Swipes"),
        ]
        main, _ = _consolidate_mailer(results)
        ids = [r.slide_id for r in main if "Jan26" in r.slide_id]
        assert ids == ["A13.Jan26", "A12.Jan26.Swipes", "A12.Jan26.Spend"]


# =============================================================================
# _result_to_slide
# =============================================================================


class TestResultToSlide:
    """Converting AnalysisResult to SlideContent."""

    def test_failed_result_returns_none(self):
        r = AnalysisResult(slide_id="DCTR-1", title="Test", success=False)
        assert _result_to_slide(r) is None

    def test_no_chart_returns_none(self):
        r = AnalysisResult(slide_id="DCTR-1", title="Test", chart_path=None)
        assert _result_to_slide(r) is None

    def test_mapped_layout_applied(self, tmp_path):
        chart = tmp_path / "chart.png"
        _write_minimal_png(chart)
        r = AnalysisResult(slide_id="A9.1", title="Attrition Rate", chart_path=chart)
        sc = _result_to_slide(r)
        assert sc is not None
        assert sc.layout_index == LAYOUT_CUSTOM
        assert sc.slide_type == "screenshot_kpi"

    def test_explicit_layout_preserved(self, tmp_path):
        chart = tmp_path / "chart.png"
        _write_minimal_png(chart)
        r = AnalysisResult(
            slide_id="A9.1",
            title="Test",
            chart_path=chart,
            layout_index=13,
            slide_type="screenshot",
        )
        sc = _result_to_slide(r)
        assert sc.layout_index == 13
        assert sc.slide_type == "screenshot"

    def test_default_uses_map(self, tmp_path):
        chart = tmp_path / "chart.png"
        _write_minimal_png(chart)
        r = AnalysisResult(slide_id="DCTR-1", title="Overall DCTR", chart_path=chart)
        sc = _result_to_slide(r)
        assert sc.layout_index == LAYOUT_CUSTOM
        assert sc.slide_type == "screenshot"


# =============================================================================
# _consolidate
# =============================================================================


class TestConsolidate:
    """Consolidation merges pairs and separates appendix."""

    def _make_result(self, slide_id, tmp_path=None):
        chart = None
        if tmp_path:
            chart = tmp_path / f"{slide_id}.png"
            _write_minimal_png(chart)
        return AnalysisResult(slide_id=slide_id, title=f"Title {slide_id}", chart_path=chart)

    def test_merge_produces_multi_screenshot(self, tmp_path):
        left = self._make_result("A7.6a", tmp_path)
        right = self._make_result("A7.4", tmp_path)
        other = self._make_result("A7.10a", tmp_path)

        main, appendix = _consolidate([left, right, other], DCTR_MERGES, DCTR_APPENDIX_IDS)
        # The merge should produce one multi_screenshot replacing both
        multi = [
            s for s in main if isinstance(s, SlideContent) and s.slide_type == "multi_screenshot"
        ]
        assert len(multi) == 1
        assert len(multi[0].images) == 2

    def test_appendix_separated(self, tmp_path):
        slides = [self._make_result(sid, tmp_path) for sid in ["A7.5", "A7.6b", "A7.10a"]]
        main, appendix = _consolidate(slides, DCTR_MERGES, DCTR_APPENDIX_IDS)
        appendix_ids = {getattr(r, "slide_id", "") for r in appendix}
        assert "A7.5" in appendix_ids
        assert "A7.6b" in appendix_ids

    def test_non_matched_preserved(self, tmp_path):
        slides = [self._make_result("A7.10a", tmp_path)]
        main, appendix = _consolidate(slides, DCTR_MERGES, DCTR_APPENDIX_IDS)
        assert len(main) == 1
        assert len(appendix) == 0

    def test_attrition_merge(self, tmp_path):
        slides = [
            self._make_result("A9.3", tmp_path),
            self._make_result("A9.6", tmp_path),
            self._make_result("A9.1", tmp_path),
        ]
        main, appendix = _consolidate(slides, ATTRITION_MERGES, ATTRITION_APPENDIX_IDS)
        multi = [s for s in main if isinstance(s, SlideContent)]
        assert len(multi) == 1  # merged pair
        assert len(main) == 2  # merged + A9.1


# =============================================================================
# _get_section / _group_by_section
# =============================================================================


class TestSectionMapping:
    """Section assignment from slide IDs."""

    @pytest.mark.parametrize(
        "slide_id,expected",
        [
            ("DCTR-1", "dctr"),
            ("A7.4", "dctr"),
            ("A8.1", "rege"),
            ("A9.1", "attrition"),
            ("A11.1", "value"),
            ("A12.Nov25.Swipes", "mailer"),
            ("A13.Sep24", "mailer"),
            ("ICS-1", "ics"),
            ("S1", "insights"),
            ("XYZ-1", "other"),
        ],
    )
    def test_section_mapping(self, slide_id, expected):
        assert _get_section(slide_id) == expected

    def test_group_counts(self):
        slides = [
            AnalysisResult(slide_id="DCTR-1", title=""),
            AnalysisResult(slide_id="DCTR-2", title=""),
            AnalysisResult(slide_id="A9.1", title=""),
        ]
        groups = _group_by_section(slides)
        assert len(groups["dctr"]) == 2
        assert len(groups["attrition"]) == 1


# =============================================================================
# _build_preamble_slides
# =============================================================================


class TestPreamble:
    """Preamble slide generation."""

    def test_produces_13_slides(self):
        slides = _build_preamble_slides("Test CU", "2026.01")
        assert len(slides) == 13

    def test_title_contains_client(self):
        slides = _build_preamble_slides("Acme Bank", "2026.02")
        assert "Acme Bank" in slides[0].title

    def test_title_contains_date(self):
        slides = _build_preamble_slides("Test CU", "2026.03")
        assert "March" in slides[0].title

    def test_all_valid_layouts(self):
        slides = _build_preamble_slides("Test CU", "2026.01")
        for i, sc in enumerate(slides):
            assert 0 <= sc.layout_index <= 19, (
                f"Preamble slide {i} has invalid layout {sc.layout_index}"
            )

    def test_master_title_uses_rpe_layout(self):
        slides = _build_preamble_slides("Test CU", "2026.01")
        assert slides[0].layout_index == LAYOUT_TITLE_RPE

    def test_slide_types(self):
        slides = _build_preamble_slides("Test CU", "2026.01")
        types = {sc.slide_type for sc in slides}
        assert "title" in types
        assert "blank" in types


# =============================================================================
# DeckBuilder class
# =============================================================================


class TestDeckBuilderClass:
    """DeckBuilder assembly with template."""

    def test_build_empty_slides(self, tmp_path):
        from ars_analysis.output.deck_builder import _FALLBACK_TEMPLATE

        if not _FALLBACK_TEMPLATE.exists():
            pytest.skip("Template not available")

        output = tmp_path / "test.pptx"
        builder = DeckBuilder(str(_FALLBACK_TEMPLATE))
        builder.build([], str(output))
        assert output.exists()
        prs = Presentation(str(output))
        assert len(prs.slides) == 0

    def test_build_section_slide(self, tmp_path):
        from ars_analysis.output.deck_builder import _FALLBACK_TEMPLATE

        if not _FALLBACK_TEMPLATE.exists():
            pytest.skip("Template not available")

        output = tmp_path / "test.pptx"
        slides = [SlideContent(slide_type="section", title="Test Section", layout_index=LAYOUT_SECTION)]
        builder = DeckBuilder(str(_FALLBACK_TEMPLATE))
        builder.build(slides, str(output))
        prs = Presentation(str(output))
        assert len(prs.slides) == 1

    def test_build_screenshot_slide(self, tmp_path):
        from ars_analysis.output.deck_builder import _FALLBACK_TEMPLATE

        if not _FALLBACK_TEMPLATE.exists():
            pytest.skip("Template not available")

        chart = tmp_path / "chart.png"
        _write_minimal_png(chart)
        output = tmp_path / "test.pptx"
        slides = [
            SlideContent(
                slide_type="screenshot",
                title="Chart Test",
                images=[str(chart)],
                layout_index=LAYOUT_CUSTOM,
            )
        ]
        builder = DeckBuilder(str(_FALLBACK_TEMPLATE))
        builder.build(slides, str(output))
        prs = Presentation(str(output))
        assert len(prs.slides) == 1


# =============================================================================
# build_deck integration
# =============================================================================


class TestBuildDeckIntegration:
    """build_deck end-to-end with minimal context."""

    def test_full_flow_with_charts(self, tmp_path):
        ctx = PipelineContext(
            client=ClientInfo(client_id="5678", client_name="Integration CU", month="2026.02"),
            paths=OutputPaths.from_base(tmp_path, "5678", "2026.02"),
        )
        ctx.paths.charts_dir.mkdir(parents=True, exist_ok=True)

        # Create charts for multiple sections
        for sid in ["DCTR-1", "A9.1", "A8.3"]:
            chart = ctx.paths.charts_dir / f"{sid}.png"
            _write_minimal_png(chart)
            ctx.all_slides.append(
                AnalysisResult(slide_id=sid, title=f"Analysis {sid}", chart_path=chart)
            )

        result = build_deck(ctx)
        assert result is not None
        assert result.exists()
        prs = Presentation(str(result))
        # 13 preamble + 3 section dividers + 3 charts + 1 summary placeholder
        assert len(prs.slides) >= 18

    def test_layout_variety(self, tmp_path):
        """Verify the deck uses multiple layout indices (not all #5)."""
        ctx = PipelineContext(
            client=ClientInfo(client_id="9999", client_name="Layout CU", month="2026.01"),
            paths=OutputPaths.from_base(tmp_path, "9999", "2026.01"),
        )
        ctx.paths.charts_dir.mkdir(parents=True, exist_ok=True)

        for sid in ["DCTR-1", "A9.1", "A7.10a"]:
            chart = ctx.paths.charts_dir / f"{sid}.png"
            _write_minimal_png(chart)
            ctx.all_slides.append(
                AnalysisResult(slide_id=sid, title=f"Analysis {sid}", chart_path=chart)
            )

        result = build_deck(ctx)
        prs = Presentation(str(result))
        layout_indices = set()
        for slide in prs.slides:
            # Get the layout index by finding matching layout
            for i, layout in enumerate(prs.slide_layouts):
                if slide.slide_layout == layout:
                    layout_indices.add(i)
                    break
        # Should have more than just 1 layout type
        assert len(layout_indices) >= 3, f"Only {layout_indices} layouts used"


# ---------------------------------------------------------------------------
# Phase 2: Section ordering + KPI dashboard
# ---------------------------------------------------------------------------


class TestSectionLabels:
    """Verify section labels are business questions."""

    def test_labels_are_questions(self):
        from ars_analysis.output.deck_builder import _SECTION_LABELS

        for key, label in _SECTION_LABELS.items():
            assert label[0].isupper(), f"{key} label doesn't start uppercase: {label}"
            assert "?" in label, f"{key} label is not a question: {label}"

    def test_section_order_matches_labels(self):
        from ars_analysis.output.deck_builder import _SECTION_LABELS, SECTION_ORDER

        for key in SECTION_ORDER:
            assert key in _SECTION_LABELS, f"SECTION_ORDER key {key!r} not in labels"

    def test_section_order_scr_arc(self):
        from ars_analysis.output.deck_builder import SECTION_ORDER

        # Overview comes first (situation)
        assert SECTION_ORDER[0] == "overview"
        # Insights comes last (resolution / call to action)
        assert SECTION_ORDER[-1] == "insights"
        # DCTR/Reg E/Attrition before Mailer (complication before resolution)
        dctr_idx = SECTION_ORDER.index("dctr")
        mailer_idx = SECTION_ORDER.index("mailer")
        assert dctr_idx < mailer_idx


class TestExecutiveKPI:
    """Verify KPI dashboard builder."""

    def test_executive_kpi_with_data(self):
        from ars_analysis.output.deck_builder import _build_executive_kpi

        ctx_results = {
            "dctr_1": {"insights": {"overall_dctr": 0.342, "total_accounts": 12400}},
            "rege_1": {"insights": {"opt_in_rate": 0.67}},
            "attrition_1": {"insights": {"overall_rate": 0.06}},
        }
        sc = _build_executive_kpi(ctx_results, title_date="January 2026")
        assert sc.slide_type == "kpi_dashboard"
        assert "Executive Dashboard" in sc.title
        assert "January 2026" in sc.title
        assert sc.kpis is not None
        assert "DCTR Penetration" in sc.kpis
        assert "34.2%" in sc.kpis["DCTR Penetration"]
        assert "green" in sc.kpis["DCTR Penetration"]  # 34.2% > 30%
        assert "Reg E Opt-In" in sc.kpis
        assert "Attrition Rate" in sc.kpis
        assert "yellow" in sc.kpis["Attrition Rate"]  # 6% is between 5-10%

    def test_executive_kpi_missing_modules(self):
        from ars_analysis.output.deck_builder import _build_executive_kpi

        sc = _build_executive_kpi({})
        assert sc.kpis is not None
        for value in sc.kpis.values():
            assert "N/A" in value or "|" not in value

    def test_executive_kpi_nan_values(self):
        from ars_analysis.output.deck_builder import _build_executive_kpi

        ctx_results = {
            "dctr_1": {"insights": {"overall_dctr": float("nan"), "total_accounts": 0}},
        }
        sc = _build_executive_kpi(ctx_results)
        assert "N/A" in sc.kpis["DCTR Penetration"]

    def test_kpi_dashboard_slide_type(self):
        """Verify DeckBuilder registers kpi_dashboard."""
        from ars_analysis.output.deck_builder import _FALLBACK_TEMPLATE

        if not _FALLBACK_TEMPLATE.exists():
            pytest.skip("Template not found")
        builder = DeckBuilder(str(_FALLBACK_TEMPLATE))
        sc = SlideContent(
            slide_type="kpi_dashboard",
            title="Test Dashboard",
            kpis={"Metric": "42%|green"},
            layout_index=LAYOUT_CUSTOM,
        )
        builder.build([sc], str(Path("/tmp/test_kpi_dashboard.pptx")))
        prs = Presentation("/tmp/test_kpi_dashboard.pptx")
        assert len(prs.slides) == 1


# ---------------------------------------------------------------------------
# Phase 3: Layout variety
# ---------------------------------------------------------------------------


class TestChartNarrativeSlide:
    """Verify chart_narrative slide type."""

    def test_chart_narrative_builds(self, tmp_path):
        from ars_analysis.output.deck_builder import _FALLBACK_TEMPLATE

        if not _FALLBACK_TEMPLATE.exists():
            pytest.skip("Template not found")

        png = tmp_path / "chart.png"
        _write_minimal_png(png)

        builder = DeckBuilder(str(_FALLBACK_TEMPLATE))
        sc = SlideContent(
            slide_type="chart_narrative",
            title="Debit adoption accelerating",
            images=[str(png)],
            bullets=["DCTR up 4pp vs prior quarter", "Branch 5 leads at 42%"],
            layout_index=LAYOUT_CUSTOM,
        )
        out = tmp_path / "test_narrative.pptx"
        builder.build([sc], str(out))
        prs = Presentation(str(out))
        assert len(prs.slides) == 1

    def test_chart_narrative_with_kpis(self, tmp_path):
        from ars_analysis.output.deck_builder import _FALLBACK_TEMPLATE

        if not _FALLBACK_TEMPLATE.exists():
            pytest.skip("Template not found")

        png = tmp_path / "chart.png"
        _write_minimal_png(png)

        builder = DeckBuilder(str(_FALLBACK_TEMPLATE))
        sc = SlideContent(
            slide_type="chart_narrative",
            title="Revenue Impact",
            images=[str(png)],
            kpis={"DCTR": "34.2%", "Revenue": "$142K"},
            layout_index=LAYOUT_CUSTOM,
        )
        out = tmp_path / "test_kpi_narrative.pptx"
        builder.build([sc], str(out))
        prs = Presentation(str(out))
        assert len(prs.slides) == 1


class TestKPIHeroSlide:
    """Verify kpi_hero slide type."""

    def test_kpi_hero_builds(self, tmp_path):
        from ars_analysis.output.deck_builder import _FALLBACK_TEMPLATE

        if not _FALLBACK_TEMPLATE.exists():
            pytest.skip("Template not found")

        builder = DeckBuilder(str(_FALLBACK_TEMPLATE))
        sc = SlideContent(
            slide_type="kpi_hero",
            title="Program Performance Snapshot",
            kpis={"DCTR": "34.2%", "Opt-In": "67%", "Attrition": "6.2%"},
            layout_index=LAYOUT_CUSTOM,
        )
        out = tmp_path / "test_hero.pptx"
        builder.build([sc], str(out))
        prs = Presentation(str(out))
        assert len(prs.slides) == 1

    def test_kpi_hero_with_chart(self, tmp_path):
        from ars_analysis.output.deck_builder import _FALLBACK_TEMPLATE

        if not _FALLBACK_TEMPLATE.exists():
            pytest.skip("Template not found")

        png = tmp_path / "chart.png"
        _write_minimal_png(png)

        builder = DeckBuilder(str(_FALLBACK_TEMPLATE))
        sc = SlideContent(
            slide_type="kpi_hero",
            title="Trend Overview",
            kpis={"DCTR": "34.2%", "Accounts": "12,400"},
            images=[str(png)],
            layout_index=LAYOUT_CUSTOM,
        )
        out = tmp_path / "test_hero_chart.pptx"
        builder.build([sc], str(out))
        prs = Presentation(str(out))
        assert len(prs.slides) == 1

    def test_kpi_hero_empty_kpis(self, tmp_path):
        from ars_analysis.output.deck_builder import _FALLBACK_TEMPLATE

        if not _FALLBACK_TEMPLATE.exists():
            pytest.skip("Template not found")

        builder = DeckBuilder(str(_FALLBACK_TEMPLATE))
        sc = SlideContent(
            slide_type="kpi_hero",
            title="Empty Dashboard",
            kpis={},
            layout_index=LAYOUT_CUSTOM,
        )
        out = tmp_path / "test_hero_empty.pptx"
        builder.build([sc], str(out))
        prs = Presentation(str(out))
        assert len(prs.slides) == 1
