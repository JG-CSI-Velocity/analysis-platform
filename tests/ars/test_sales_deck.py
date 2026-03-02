"""Tests for RPE sales conference deck builder."""

from __future__ import annotations

from pathlib import Path

import pytest


class TestSalesDeckConfig:
    def test_default_config(self):
        from ars_analysis.output.sales_deck_builder import SalesDeckConfig

        cfg = SalesDeckConfig()
        assert cfg.conference_name == "2026 Credit Union Conference"
        assert cfg.tagline == "The Complete Debit Card Lifecycle Platform"
        assert cfg.mrpc_chart_path is None
        assert cfg.output_path == Path("RPE_Sales_Deck.pptx")

    def test_custom_config(self):
        from ars_analysis.output.sales_deck_builder import SalesDeckConfig

        cfg = SalesDeckConfig(
            conference_name="CUNA GAC 2026",
            booth_number="123",
            contact_info="sales@csi.com",
        )
        assert cfg.conference_name == "CUNA GAC 2026"
        assert cfg.booth_number == "123"

    def test_config_is_frozen(self):
        from ars_analysis.output.sales_deck_builder import SalesDeckConfig

        cfg = SalesDeckConfig()
        with pytest.raises(AttributeError):
            cfg.conference_name = "Changed"


class TestSalesCharts:
    def test_lifecycle_diagram(self, tmp_path: Path):
        from ars_analysis.output.sales_charts import lifecycle_diagram

        result = lifecycle_diagram(tmp_path)
        assert result.exists()
        assert result.suffix == ".png"
        assert result.stat().st_size > 1000

    def test_ics_source_chart(self, tmp_path: Path):
        from ars_analysis.output.sales_charts import ics_source_chart

        result = ics_source_chart(tmp_path)
        assert result.exists()
        assert result.suffix == ".png"

    def test_service_adoption_chart(self, tmp_path: Path):
        from ars_analysis.output.sales_charts import service_adoption_chart

        result = service_adoption_chart(tmp_path)
        assert result.exists()
        assert result.suffix == ".png"

    def test_swipe_ladder_chart(self, tmp_path: Path):
        from ars_analysis.output.sales_charts import swipe_ladder_chart

        result = swipe_ladder_chart(tmp_path)
        assert result.exists()
        assert result.suffix == ".png"

    def test_competition_chart(self, tmp_path: Path):
        from ars_analysis.output.sales_charts import competition_chart

        result = competition_chart(tmp_path)
        assert result.exists()
        assert result.suffix == ".png"

    def test_financial_services_chart(self, tmp_path: Path):
        from ars_analysis.output.sales_charts import financial_services_chart

        result = financial_services_chart(tmp_path)
        assert result.exists()
        assert result.suffix == ".png"

    def test_lifecycle_kpi_dashboard(self, tmp_path: Path):
        from ars_analysis.output.sales_charts import lifecycle_kpi_dashboard

        result = lifecycle_kpi_dashboard(tmp_path)
        assert result.exists()
        assert result.suffix == ".png"

    def test_mrpc_fallback_chart(self, tmp_path: Path):
        from ars_analysis.output.sales_charts import mrpc_fallback_chart

        result = mrpc_fallback_chart(tmp_path)
        assert result.exists()
        assert result.suffix == ".png"


class TestSlideDefinitions:
    def test_slide_count(self):
        from ars_analysis.output.sales_deck_builder import SalesDeckConfig, _build_slide_definitions

        cfg = SalesDeckConfig()
        # Provide dummy chart paths (don't need to exist for definition building)
        charts = {k: "/tmp/fake.png" for k in [
            "lifecycle", "ics_source", "service_adoption", "swipe_ladder",
            "competition", "financial_services", "lifecycle_kpi", "mrpc",
        ]}
        slides = _build_slide_definitions(cfg, charts)
        assert len(slides) == 13

    def test_first_slide_is_title(self):
        from ars_analysis.output.sales_deck_builder import SalesDeckConfig, _build_slide_definitions

        cfg = SalesDeckConfig()
        charts = {k: "/tmp/fake.png" for k in [
            "lifecycle", "ics_source", "service_adoption", "swipe_ladder",
            "competition", "financial_services", "lifecycle_kpi", "mrpc",
        ]}
        slides = _build_slide_definitions(cfg, charts)
        assert slides[0].slide_type == "title"
        assert "Retail Performance Engine" in slides[0].title

    def test_last_slide_is_cta(self):
        from ars_analysis.output.sales_deck_builder import SalesDeckConfig, _build_slide_definitions

        cfg = SalesDeckConfig(booth_number="456")
        charts = {k: "/tmp/fake.png" for k in [
            "lifecycle", "ics_source", "service_adoption", "swipe_ladder",
            "competition", "financial_services", "lifecycle_kpi", "mrpc",
        ]}
        slides = _build_slide_definitions(cfg, charts)
        assert "Booth #456" in slides[-1].title

    def test_all_slides_have_notes(self):
        from ars_analysis.output.sales_deck_builder import SalesDeckConfig, _build_slide_definitions

        cfg = SalesDeckConfig()
        charts = {k: "/tmp/fake.png" for k in [
            "lifecycle", "ics_source", "service_adoption", "swipe_ladder",
            "competition", "financial_services", "lifecycle_kpi", "mrpc",
        ]}
        slides = _build_slide_definitions(cfg, charts)
        for i, slide in enumerate(slides):
            assert slide.notes_text, f"Slide {i + 1} missing speaker notes"

    def test_kpi_hero_slide(self):
        from ars_analysis.output.sales_deck_builder import SalesDeckConfig, _build_slide_definitions

        cfg = SalesDeckConfig()
        charts = {k: "/tmp/fake.png" for k in [
            "lifecycle", "ics_source", "service_adoption", "swipe_ladder",
            "competition", "financial_services", "lifecycle_kpi", "mrpc",
        ]}
        slides = _build_slide_definitions(cfg, charts)
        kpi_slide = slides[2]
        assert kpi_slide.slide_type == "kpi_hero"
        assert len(kpi_slide.kpis) == 4


class TestBuildSalesDeck:
    def test_generates_pptx(self, tmp_path: Path):
        from ars_analysis.output.sales_deck_builder import SalesDeckConfig, build_sales_deck

        output = tmp_path / "test_deck.pptx"
        cfg = SalesDeckConfig(output_path=output)
        result = build_sales_deck(cfg)
        assert result.exists()
        assert result.suffix == ".pptx"
        assert result.stat().st_size > 10000

    def test_slide_count_in_pptx(self, tmp_path: Path):
        from pptx import Presentation

        from ars_analysis.output.sales_deck_builder import SalesDeckConfig, build_sales_deck

        output = tmp_path / "deck.pptx"
        cfg = SalesDeckConfig(output_path=output)
        build_sales_deck(cfg)

        prs = Presentation(str(output))
        assert len(prs.slides) == 13

    def test_with_mrpc_chart(self, tmp_path: Path):
        from ars_analysis.output.sales_charts import mrpc_fallback_chart
        from ars_analysis.output.sales_deck_builder import SalesDeckConfig, build_sales_deck

        # Generate a chart to use as "user-provided"
        mrpc_img = mrpc_fallback_chart(tmp_path)

        output = tmp_path / "deck_mrpc.pptx"
        cfg = SalesDeckConfig(
            mrpc_chart_path=mrpc_img,
            output_path=output,
        )
        result = build_sales_deck(cfg)
        assert result.exists()

    def test_with_booth_and_contact(self, tmp_path: Path):
        from pptx import Presentation

        from ars_analysis.output.sales_deck_builder import SalesDeckConfig, build_sales_deck

        output = tmp_path / "deck_booth.pptx"
        cfg = SalesDeckConfig(
            booth_number="789",
            contact_info="sales@csi.com",
            output_path=output,
        )
        build_sales_deck(cfg)

        prs = Presentation(str(output))
        # Last slide should have booth number
        last_slide = prs.slides[-1]
        texts = []
        for shape in last_slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    texts.append(p.text)
        all_text = " ".join(texts)
        assert "789" in all_text or "Booth" in all_text
