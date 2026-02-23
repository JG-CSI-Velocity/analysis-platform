"""Tests for Primary/Secondary deck consolidation."""

import struct
import zlib

import pandas as pd
import pytest
from pptx import Presentation

from ics_toolkit.analysis.analyses.base import AnalysisResult
from ics_toolkit.analysis.exports.pptx import (
    _PRIMARY_NAMES,
    KPI_PANEL_ANALYSES,
    PRIMARY_STORYLINE,
    SECTION_MAP,
    write_ics_reports,
    write_pptx_primary,
    write_pptx_secondary,
)


def _make_tiny_png() -> bytes:
    """Create a minimal valid 1x1 PNG for testing."""
    raw = b"\x00\xff\x00\x00"
    compressed = zlib.compress(raw)
    ihdr = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)

    def chunk(tag, data):
        c = tag + data
        return struct.pack(">I", len(data)) + c + struct.pack(">I", zlib.crc32(c) & 0xFFFFFFFF)

    return (
        b"\x89PNG\r\n\x1a\n"
        + chunk(b"IHDR", ihdr)
        + chunk(b"IDAT", compressed)
        + chunk(b"IEND", b"")
    )


def _make_analysis(name: str, title: str | None = None, metadata: dict | None = None):
    """Create a mock AnalysisResult with a 2-row DataFrame."""
    return AnalysisResult.from_df(
        name,
        title or name,
        pd.DataFrame({"Metric": ["A", "B"], "Value": [10, 20]}),
        sheet_name=name.replace(" ", "_")[:20],
        metadata=metadata,
    )


@pytest.fixture
def all_analyses():
    """Create mock analyses for every name in SECTION_MAP."""
    analyses = []
    for names in SECTION_MAP.values():
        for name in names:
            meta = {}
            if name in KPI_PANEL_ANALYSES:
                meta = {
                    "hero_kpis": {
                        "Total Accounts": 1500,
                        "Penetration Rate": 32.5,
                        "Active %": 68.0,
                        "Avg Interchange": 12.50,
                    },
                }
            analyses.append(_make_analysis(name, metadata=meta))
    return analyses


@pytest.fixture
def all_chart_pngs():
    """Create a PNG for every analysis that would have a chart."""
    png = _make_tiny_png()
    pngs = {}
    for names in SECTION_MAP.values():
        for name in names:
            pngs[name] = png
    return pngs


class TestPrimaryStorylineConfig:
    def test_has_six_sections(self):
        assert len(PRIMARY_STORYLINE) == 6

    def test_section_titles_are_questions(self):
        for title in PRIMARY_STORYLINE:
            assert title.endswith("?"), f"Section title should be a question: {title}"

    def test_primary_names_subset_of_section_map(self):
        all_section_names = set()
        for names in SECTION_MAP.values():
            all_section_names.update(names)
        orphans = _PRIMARY_NAMES - all_section_names
        assert not orphans, f"Primary references analyses not in SECTION_MAP: {orphans}"

    def test_kpi_panel_analyses_subset_of_primary(self):
        assert KPI_PANEL_ANALYSES.issubset(_PRIMARY_NAMES), (
            f"KPI panels not in primary: {KPI_PANEL_ANALYSES - _PRIMARY_NAMES}"
        )


class TestPrimaryDeck:
    def test_creates_file(self, sample_settings, all_analyses, tmp_path):
        path = tmp_path / "primary.pptx"
        result = write_pptx_primary(
            sample_settings, all_analyses, output_path=path, chart_pngs={},
        )
        assert result.exists()

    def test_slide_count_under_35(self, sample_settings, all_analyses, all_chart_pngs, tmp_path):
        path = tmp_path / "primary.pptx"
        write_pptx_primary(
            sample_settings, all_analyses, output_path=path, chart_pngs=all_chart_pngs,
        )
        prs = Presentation(str(path))
        count = len(prs.slides)
        assert count <= 35, f"Primary deck has {count} slides, expected <= 35"
        assert count >= 5, f"Primary deck has only {count} slides, too few"

    def test_storyline_sections_present(
        self, sample_settings, all_analyses, all_chart_pngs, tmp_path,
    ):
        """All 6 storyline sections should have divider slides."""
        path = tmp_path / "primary.pptx"
        write_pptx_primary(
            sample_settings, all_analyses, output_path=path, chart_pngs=all_chart_pngs,
        )
        prs = Presentation(str(path))
        all_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    all_text.append(shape.text_frame.text)

        combined = " ".join(all_text)
        for section_title in PRIMARY_STORYLINE:
            assert section_title in combined, f"Section '{section_title}' not found in deck"

    def test_merged_slides_have_two_images(
        self, sample_settings, all_analyses, all_chart_pngs, tmp_path,
    ):
        """Merge-pair slides should contain exactly 2 picture shapes."""
        path = tmp_path / "primary.pptx"
        write_pptx_primary(
            sample_settings, all_analyses, output_path=path, chart_pngs=all_chart_pngs,
        )
        prs = Presentation(str(path))

        merged_slides = []
        for slide in prs.slides:
            texts = [
                s.text_frame.text for s in slide.shapes if s.has_text_frame
            ]
            combined = " ".join(texts)
            if "|" in combined:
                merged_slides.append(slide)

        assert len(merged_slides) > 0, "No merged slides found"
        for slide in merged_slides:
            pics = [s for s in slide.shapes if s.shape_type == 13]  # 13 = Picture
            assert len(pics) == 2, (
                f"Merged slide should have 2 pictures, got {len(pics)}"
            )

    def test_no_table_in_primary(
        self, sample_settings, all_analyses, all_chart_pngs, tmp_path,
    ):
        """Primary deck should not contain data tables (tables are in secondary)."""
        path = tmp_path / "primary.pptx"
        write_pptx_primary(
            sample_settings, all_analyses, output_path=path, chart_pngs=all_chart_pngs,
        )
        prs = Presentation(str(path))
        table_count = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_table:
                    table_count += 1
        assert table_count == 0, f"Primary deck has {table_count} tables, expected 0"

    def test_kpi_panel_for_summary(self, sample_settings, tmp_path):
        """KPI analyses without charts should produce KPI panel slides."""
        analyses = [
            _make_analysis(
                "Executive Summary",
                metadata={
                    "hero_kpis": {"Total": 1500, "Rate": 32.5, "Active": 68},
                },
            ),
        ]
        path = tmp_path / "kpi.pptx"
        write_pptx_primary(
            sample_settings, analyses, output_path=path, chart_pngs={},
        )
        prs = Presentation(str(path))
        # Title + section divider + KPI slide = at least 3
        assert len(prs.slides) >= 3


class TestSecondaryDeck:
    def test_creates_file(self, sample_settings, all_analyses, tmp_path):
        path = tmp_path / "secondary.pptx"
        result = write_pptx_secondary(sample_settings, all_analyses, output_path=path)
        assert result.exists()

    def test_all_analyses_present(self, sample_settings, all_analyses, tmp_path):
        """Every successful analysis name should appear in secondary deck text."""
        path = tmp_path / "secondary.pptx"
        write_pptx_secondary(sample_settings, all_analyses, output_path=path)
        prs = Presentation(str(path))
        all_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    all_text.append(shape.text_frame.text)
        combined = " ".join(all_text)

        successful = [a for a in all_analyses if a.error is None and not a.df.empty]
        for a in successful:
            assert a.title in combined, f"Analysis '{a.title}' missing from secondary deck"

    def test_has_appendix_title(self, sample_settings, all_analyses, tmp_path):
        path = tmp_path / "secondary.pptx"
        write_pptx_secondary(sample_settings, all_analyses, output_path=path)
        prs = Presentation(str(path))
        first_slide = prs.slides[0]
        texts = [s.text_frame.text for s in first_slide.shapes if s.has_text_frame]
        combined = " ".join(texts)
        assert "Detail Appendix" in combined


class TestBothDecks:
    def test_write_ics_reports_creates_both(
        self, sample_settings, all_analyses, all_chart_pngs, tmp_path,
    ):
        primary, secondary = write_ics_reports(
            sample_settings,
            all_analyses,
            chart_pngs=all_chart_pngs,
            output_dir=tmp_path,
        )
        assert primary.exists()
        assert secondary.exists()
        assert "Primary" in primary.name
        assert "Secondary" in secondary.name

    def test_primary_much_smaller_than_secondary(
        self, sample_settings, all_analyses, all_chart_pngs, tmp_path,
    ):
        primary_path, secondary_path = write_ics_reports(
            sample_settings,
            all_analyses,
            chart_pngs=all_chart_pngs,
            output_dir=tmp_path,
        )
        primary_prs = Presentation(str(primary_path))
        secondary_prs = Presentation(str(secondary_path))
        assert len(primary_prs.slides) < len(secondary_prs.slides)

    def test_no_analysis_lost(
        self, sample_settings, all_analyses, all_chart_pngs, tmp_path,
    ):
        """Every successful analysis should appear in at least one deck."""
        primary_path, secondary_path = write_ics_reports(
            sample_settings,
            all_analyses,
            chart_pngs=all_chart_pngs,
            output_dir=tmp_path,
        )

        # All successful analyses appear in the secondary deck (which has everything)
        secondary_prs = Presentation(str(secondary_path))
        all_text = []
        for slide in secondary_prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    all_text.append(shape.text_frame.text)
        combined = " ".join(all_text)

        successful = [a for a in all_analyses if a.error is None and not a.df.empty]
        for a in successful:
            assert a.title in combined, f"Analysis '{a.title}' lost from both decks"
